import os
import io
import datetime
import json
import time
import requests
import streamlit as st
import sqlite3
import re
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FutureTimeoutError
from dotenv import load_dotenv
from langchain_google_genai import ChatGoogleGenerativeAI
try:
    from langchain_core.prompts import ChatPromptTemplate
except ImportError:
    from langchain.prompts import ChatPromptTemplate
from langchain.prompts import PromptTemplate
from langchain.chains.summarize import load_summarize_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from langchain_core.output_parsers import StrOutputParser
from langchain_google_community import GoogleSearchAPIWrapper

# --- 0. 인증/회원관리 ---
from auth import require_login, render_sidebar_user_panel, current_user, is_admin

# --- 0-1. 백엔드 콘솔 로깅 ---
from logging_setup import setup_logging, get_logger, log_event
setup_logging()
log = get_logger("APP")

# --- 0-2. 공통 UI 테마 / API 상태 ---
from ui_theme import inject_global_css, page_header, render_stepper
from api_status import render_api_status_bar
from ppt_llm import get_ppt_llm
# Streamlit은 매 인터랙션마다 스크립트를 재실행하므로 1회만 찍는 가드
if not getattr(setup_logging, "_boot_logged", False):
    log.info("==== 서비스 로딩 시작 (Proposal Generator) ====")
    log.info(f"Python={os.sys.version.split()[0]} CWD={os.getcwd()}")
    setup_logging._boot_logged = True

# --- 1. 환경변수 및 페이지 설정 ---
load_dotenv()
gemini_api_key = os.getenv("GEMINI_API_KEY") or st.secrets.get("GEMINI_API_KEY", "")
google_api_key = os.getenv("GOOGLE_API_KEY") or st.secrets.get("GOOGLE_API_KEY", "")
google_cse_id = os.getenv("GOOGLE_CSE_ID") or st.secrets.get("GOOGLE_CSE_ID", "")

os.environ["GOOGLE_API_KEY"] = google_api_key
os.environ["GOOGLE_CSE_ID"] = google_cse_id

def _parse_retry_timeouts(raw_value, default_sequence):
    try:
        parsed = [int(item.strip()) for item in raw_value.split(",") if item.strip()]
        return parsed or default_sequence
    except ValueError:
        return default_sequence

LLM_RETRY_TIMEOUTS = _parse_retry_timeouts(os.getenv("LLM_RETRY_TIMEOUTS", "60,120"), [60, 120])
LLM_TIMEOUT_SECONDS = max(LLM_RETRY_TIMEOUTS)
# Gemini SDK(langchain_google_genai) 내부 tenacity 재시도 횟수.
# 기본값 6은 너무 커서 504(DeadlineExceeded)가 반복되면 한 시도가 수 분씩 늘어진다.
# 우리 execute_with_retries가 시도 자체를 한번 더 감싸므로 SDK는 1회만 시도하도록 낮춘다.
LLM_SDK_MAX_RETRIES = int(os.getenv("LLM_SDK_MAX_RETRIES", "1"))
ENHANCEMENT_LOG_LIMIT = 12

st.set_page_config(page_title="제안서 & 추진계획서 자동 생성 Agent", layout="wide")
inject_global_css()
# --- 2. 데이터베이스 관리 함수 ---
DB_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "history.db")

def _connect_db():
    """모든 DB 호출에서 공통으로 사용할 커넥션 (락 대기·자동 커밋·UTF-8 친화)."""
    conn = sqlite3.connect(DB_FILE, timeout=10)
    return conn

def init_db():
    conn = _connect_db()
    c = conn.cursor()
    # WAL 모드: 강화 작업 중 History 페이지 등 동시 조회 시 락 충돌을 줄인다.
    try:
        c.execute("PRAGMA journal_mode=WAL")
        c.execute("PRAGMA synchronous=NORMAL")
    except Exception:
        pass
    c.execute('''
        CREATE TABLE IF NOT EXISTS projects (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            topic TEXT,
            timestamp TEXT NOT NULL,
            page_count INTEGER DEFAULT 15
        )
    ''')
    c.execute("PRAGMA table_info(projects)")
    columns = [row[1] for row in c.fetchall()]
    if 'page_count' not in columns:
        c.execute("ALTER TABLE projects ADD COLUMN page_count INTEGER DEFAULT 15")

    c.execute('''
        CREATE TABLE IF NOT EXISTS project_stages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            project_id INTEGER NOT NULL,
            stage_name TEXT NOT NULL,
            content TEXT,
            llm_type TEXT,
            UNIQUE(project_id, stage_name)
        )
    ''')
    # 중단 시점 재사용을 위한 섹션 단위 진행분 저장 테이블 (P1-2)
    c.execute('''
        CREATE TABLE IF NOT EXISTS stage_progress (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            project_id INTEGER NOT NULL,
            stage TEXT NOT NULL,          -- 'stage3_body' | 'stage4_enhance_body' | 'stage4_enhance_toc'
            section_key TEXT NOT NULL,    -- original_line 또는 heading 문자열
            payload TEXT,                 -- 처리 결과 본문/제목
            status TEXT,                  -- 'ok' | 'fallback_original' | 'skipped'
            attempt INTEGER DEFAULT 1,
            elapsed_sec REAL DEFAULT 0,
            updated_at TEXT DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(project_id, stage, section_key)
        )
    ''')
    conn.commit()
    conn.close()


# --- 진행분(중간 산출물) 영속화 헬퍼 ---
def save_progress_item(project_id, stage, section_key, payload, status="ok", attempt=1, elapsed_sec=0.0):
    """섹션 1개 처리 직후 즉시 저장. 같은 키면 덮어쓰기."""
    if not project_id:
        return
    try:
        conn = _connect_db()
        c = conn.cursor()
        c.execute('''
            INSERT INTO stage_progress (project_id, stage, section_key, payload, status, attempt, elapsed_sec, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, CURRENT_TIMESTAMP)
            ON CONFLICT(project_id, stage, section_key) DO UPDATE SET
                payload=excluded.payload,
                status=excluded.status,
                attempt=excluded.attempt,
                elapsed_sec=excluded.elapsed_sec,
                updated_at=CURRENT_TIMESTAMP
        ''', (project_id, stage, section_key, payload, status, attempt, elapsed_sec))
        conn.commit()
        conn.close()
    except Exception as e:
        # 영속화 실패가 작업 자체를 막지 않게 한다.
        print(f"[stage_progress 저장 실패] {e}")

def load_progress_map(project_id, stage):
    """{section_key: (payload, status)} 형태로 반환."""
    if not project_id:
        return {}
    try:
        conn = _connect_db()
        c = conn.cursor()
        c.execute('''
            SELECT section_key, payload, status, updated_at
            FROM stage_progress
            WHERE project_id = ? AND stage = ?
        ''', (project_id, stage))
        rows = c.fetchall()
        conn.close()
        return {r[0]: {"payload": r[1], "status": r[2], "updated_at": r[3]} for r in rows}
    except Exception as e:
        print(f"[stage_progress 조회 실패] {e}")
        return {}

def clear_progress(project_id, stage=None):
    """특정 stage 또는 프로젝트 전체 진행분 삭제."""
    if not project_id:
        return
    try:
        conn = _connect_db()
        c = conn.cursor()
        if stage:
            c.execute("DELETE FROM stage_progress WHERE project_id = ? AND stage = ?", (project_id, stage))
        else:
            c.execute("DELETE FROM stage_progress WHERE project_id = ?", (project_id,))
        conn.commit()
        conn.close()
    except Exception as e:
        print(f"[stage_progress 삭제 실패] {e}")

def get_progress_summary(project_id, stage):
    """재개 카드 UI에 표시할 요약 정보."""
    items = load_progress_map(project_id, stage)
    if not items:
        return None
    latest = max((v.get("updated_at") or "") for v in items.values())
    fallbacks = [k for k, v in items.items() if (v.get("status") == "fallback_original")]
    return {"count": len(items), "latest": latest, "fallback_count": len(fallbacks), "fallback_keys": fallbacks}


def verify_progress_integrity(project_id, stage, expected_keys):
    """
    재실행/이어서 실행 직전에 누락·중복·고아 항목을 점검한다.
    - missing : expected에 있지만 stage_progress에 없음
    - duplicate : (UNIQUE 제약으로 거의 없으나) 같은 section_key가 2회 이상 저장된 경우
    - orphan : DB엔 있지만 현재 expected에 없는 키 (목차가 바뀐 경우)
    - fallback : status='fallback_original'로 저장된 항목 (재시도 권장)
    """
    items = load_progress_map(project_id, stage)
    expected_set = {str(k).strip() for k in (expected_keys or []) if str(k).strip()}
    cached_set = set(items.keys())
    missing = sorted(expected_set - cached_set)
    orphan = sorted(cached_set - expected_set)
    fallback = sorted([k for k, v in items.items() if v.get("status") == "fallback_original"])
    # 중복 점검 (UNIQUE 제약이 있더라도 방어적으로 COUNT)
    duplicates = []
    try:
        conn = _connect_db()
        c = conn.cursor()
        c.execute(
            "SELECT section_key, COUNT(*) AS n FROM stage_progress "
            "WHERE project_id = ? AND stage = ? GROUP BY section_key HAVING n > 1",
            (project_id, stage),
        )
        duplicates = [row[0] for row in c.fetchall()]
        conn.close()
    except Exception as e:
        print(f"[verify_progress_integrity 중복 점검 실패] {e}")
    report = {
        "expected": len(expected_set),
        "cached": len(cached_set),
        "missing": missing,
        "orphan": orphan,
        "fallback": fallback,
        "duplicates": duplicates,
    }
    # --- 백엔드 콘솔 로깅 ---
    _lg = get_logger("VERIFY")
    _lg.info(
        f"stage={stage} project={project_id} expected={report['expected']} cached={report['cached']} "
        f"missing={len(missing)} duplicate={len(duplicates)} orphan={len(orphan)} fallback={len(fallback)}"
    )
    if missing:
        _lg.warning(f"누락 섹션 {len(missing)}개: {missing[:5]}{'...' if len(missing) > 5 else ''}")
    if duplicates:
        _lg.error(f"중복 섹션 {len(duplicates)}개: {duplicates}")
    if fallback:
        _lg.warning(f"폴백 섹션 {len(fallback)}개: {fallback[:5]}{'...' if len(fallback) > 5 else ''}")
    return report


def render_integrity_report(report, title="진행분 검증 결과"):
    """검증 리포트를 사용자에게 보여 주는 공용 UI."""
    if not report:
        return
    cols = st.columns(4)
    cols[0].metric("예상 섹션", report["expected"])
    cols[1].metric("저장된 섹션", report["cached"])
    cols[2].metric("누락", len(report["missing"]), delta_color="inverse")
    cols[3].metric("폴백/중복/고아", len(report["fallback"]) + len(report["duplicates"]) + len(report["orphan"]),
                   delta_color="inverse")
    if report["missing"]:
        with st.expander(f"⛔ 누락 섹션 {len(report['missing'])}개 (재실행 시 새로 생성됨)", expanded=True):
            for k in report["missing"]:
                st.caption(f"• {k}")
    if report["fallback"]:
        with st.expander(f"⚠️ 폴백(원문 유지) 섹션 {len(report['fallback'])}개 — 다시 시도하면 재생성됩니다", expanded=False):
            for k in report["fallback"]:
                st.caption(f"• {k}")
    if report["duplicates"]:
        st.error(f"🚨 중복 저장된 섹션 {len(report['duplicates'])}개 발견: {', '.join(report['duplicates'])}")
    if report["orphan"]:
        with st.expander(f"🧹 현재 목차에 없는 잔여 항목 {len(report['orphan'])}개 (목차 변경 흔적)", expanded=False):
            for k in report["orphan"]:
                st.caption(f"• {k}")


def notify(level, message, toast=True):
    """단계 진행 중 발생하는 알림을 통일된 방식으로 노출. level: 'info'|'success'|'warning'|'error'"""
    # 콘솔 로깅 (서버 운영자가 백엔드에서 확인)
    _lg = get_logger("NOTIFY")
    _level_map = {"info": "info", "success": "info", "warning": "warning", "error": "error"}
    getattr(_lg, _level_map.get(level, "info"))(f"[{level.upper()}] {message}")
    fn = {"info": st.info, "success": st.success, "warning": st.warning, "error": st.error}.get(level, st.info)
    try:
        fn(message)
    except Exception:
        print(f"[notify-{level}] {message}")
    if toast:
        icon = {"info": "ℹ️", "success": "✅", "warning": "⚠️", "error": "🚨"}.get(level, "ℹ️")
        try:
            st.toast(message, icon=icon)
        except Exception:
            pass

def create_new_project():
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    c.execute("INSERT INTO projects (topic, timestamp, page_count) VALUES (?, ?, ?)", ("(주제 미정)", timestamp, 15))
    project_id = c.lastrowid
    conn.commit()
    conn.close()
    return project_id

def update_project_topic(project_id, topic):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("UPDATE projects SET topic = ? WHERE id = ?", (topic, project_id))
    conn.commit()
    conn.close()

def update_project_page_count(project_id, page_count):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("UPDATE projects SET page_count = ? WHERE id = ?", (page_count, project_id))
    conn.commit()
    conn.close()

def save_stage_result(project_id, stage_name, content, llm_type=None):
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("SELECT id FROM project_stages WHERE project_id = ? AND stage_name = ?", (project_id, stage_name))
    result = c.fetchone()
    
    if result:
        stage_id = result[0]
        c.execute('''
            UPDATE project_stages 
            SET content = ?, llm_type = ?
            WHERE id = ?
        ''', (content, llm_type, stage_id))
    else:
        c.execute('''
            INSERT INTO project_stages (project_id, stage_name, content, llm_type)
            VALUES (?, ?, ?, ?)
        ''', (project_id, stage_name, content, llm_type))
        
    conn.commit()
    conn.close()

def get_project_data(project_id):
    conn = sqlite3.connect(DB_FILE)
    conn.row_factory = sqlite3.Row
    c = conn.cursor()
    c.execute("SELECT * FROM projects WHERE id = ?", (project_id,))
    project_data = c.fetchone()
    
    c.execute("SELECT stage_name, content FROM project_stages WHERE project_id = ?", (project_id,))
    stages = {row['stage_name']: row['content'] for row in c.fetchall()}
    conn.close()
    return dict(project_data) if project_data else None, stages

def load_project_into_session(project_id):
    project_data, stages = get_project_data(project_id)
    if not project_data: return

    st.session_state.project_id = project_id
    st.session_state.page_count = project_data.get('page_count', 15)
    
    if '1단계: 자료 업로드' in stages:
        try:
            docs_content = json.loads(stages['1단계: 자료 업로드'])
            st.session_state.docs = [type('Doc', (), {'page_content': c})() for c in docs_content]
        except (json.JSONDecodeError, TypeError):
            st.session_state.docs = []
    
    if '2단계: 주제 확정' in stages:
        st.session_state.finalized_topic = stages['2단계: 주제 확정']
        
    if '2단계: 목차 확정' in stages:
        st.session_state.finalized_toc = stages['2단계: 목차 확정']
        st.session_state.editable_toc = stages['2단계: 목차 확정']

    if '3단계: 본문 생성' in stages:
        st.session_state.draft_proposal = stages['3단계: 본문 생성']
    
    if '4단계: 최종본' in stages:
        st.session_state.final_proposal = stages['4단계: 최종본']


# --- 3. 백엔드 함수 정의 ---
def validate_and_normalize_toc(toc_text):
    """
    목차 텍스트를 검증하고 정규화합니다.
    - 번호 형식 통일 (1. / 1.1. / 1.1.1. 등)
    - 이모지·글머리기호 제거
    - 들여쓰기 제거
    - 번호 없는 행 필터링
    반환: (정규화된 목차 str, 경고 메시지 list)
    """
    warnings_list = []
    normalized_lines = []
    number_pattern = re.compile(r'^(\d+(?:\.\d+)*\.?)\s+(.+)')
    bullet_pattern = re.compile(r'^[-*•]\s+')
    emoji_pattern = re.compile(
        r'[\U00010000-\U0010ffff'
        r'\U0001F300-\U0001F9FF'
        r'\U00002600-\U000027BF'
        r'\U0000FE00-\U0000FE0F]+',
        flags=re.UNICODE
    )

    for raw_line in toc_text.split('\n'):
        line = raw_line.strip()
        if not line:
            continue

        # 글머리기호 제거
        if bullet_pattern.match(line):
            line = bullet_pattern.sub('', line).strip()
            warnings_list.append(f"글머리기호 제거됨: '{line}'")

        # 이모지 제거
        cleaned = emoji_pattern.sub('', line).strip()
        if cleaned != line:
            warnings_list.append(f"이모지 제거됨: '{line}' → '{cleaned}'")
            line = cleaned

        # 번호 형식 검증
        m = number_pattern.match(line)
        if not m:
            # 번호 없는 항목 → 건너뜀
            warnings_list.append(f"번호 없는 항목 제외됨: '{line}'")
            continue

        number_part = m.group(1)
        title_part = m.group(2).strip()

        # 번호 끝에 마침표 통일 (없으면 추가)
        if not number_part.endswith('.'):
            number_part += '.'
            warnings_list.append(f"번호 마침표 추가됨: '{line}'")

        normalized_lines.append(f"{number_part} {title_part}")

    if not normalized_lines:
        warnings_list.append("⚠️ 유효한 목차 항목이 없습니다. 목차를 다시 확인하세요.")

    return '\n'.join(normalized_lines), warnings_list


def get_heading_level_from_number(number_str):
    """번호 문자열에서 정확한 헤딩 레벨을 계산합니다. (점 카운팅 대신 사용)
    예: '1.' → 1, '1.1.' → 2, '1.1.1.' → 3
    """
    parts = [p for p in number_str.rstrip('.').split('.') if p]
    return min(len(parts), 3)


def get_search_tool():
    if not (google_api_key and google_cse_id):
        st.error("Google API 키 또는 검색엔진 ID가 설정되지 않았습니다. .env 파일을 확인해주세요.")
        return None
    try:
        return GoogleSearchAPIWrapper()
    except Exception as e:
        st.error(f"Google 검색 도구 초기화 중 오류가 발생했습니다: {e}")
        return None

def get_title_recommendations(documents):
    # 단일 호출 방식: 문서 전체를 15000자로 잘라 한 번에 전송 (map_reduce 대비 수십 배 빠름)
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.5, google_api_key=gemini_api_key, max_retries=LLM_SDK_MAX_RETRIES)
    full_text = "\n\n".join([doc.page_content for doc in documents])[:15000]
    prompt = PromptTemplate.from_template(
        "다음 문서 내용을 분석하여, 전체 내용을 가장 잘 대표하는 전문적인 '제안서 주제명' 5가지를 추천해주세요.\n"
        "각 주제명은 간결해야 하며, 번호를 붙여 목록으로만 답하세요. 설명 없이 주제명만 작성하세요.\n\n"
        "### 문서 내용:\n{text}\n\n### 추천 주제명 (5가지):"
    )
    chain = prompt | llm | StrOutputParser()
    recommendations_str = chain.invoke({"text": full_text})
    return [line.strip().split('. ', 1)[-1] for line in recommendations_str.split('\n') if line.strip()][:5]

def generate_detailed_toc(topic, core_toc, documents):
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-pro", temperature=0.3, google_api_key=gemini_api_key, max_retries=LLM_SDK_MAX_RETRIES)
    full_text = "\n\n".join([doc.page_content for doc in documents])
    prompt = PromptTemplate.from_template(
        "당신은 전문 제안서 컨설턴트입니다. 아래 '확정된 주제'를 명확히 인지하고, 주어진 '핵심 목차'의 구조는 유지하면서, '참고 자료'를 분석하여 가장 관련성 높은 내용으로 전문적인 '세부 목차'를 제안해주세요.\n"
        "**[규칙]**\n"
        "1. 번호는 '1.1.', '1.2.', '2.1.'과 같은 계층적 형식으로 붙여주세요.\n"
        "2. **[매우 중요]** `-`나 `*` 같은 글머리 기호는 절대 사용하지 마세요.\n"
        "3. **[절대 금지]** 목차 항목에 📌, 🎯 와 같은 이모지(Emoji)나 특수기호는 절대 사용하지 마세요.\n\n"
        "### 확정된 주제:\n{topic}\n\n"
        "### 핵심 목차 (구조 유지):\n{core_toc}\n\n"
        "### 참고 자료:\n{context}\n\n"
        "### 생성할 상세 목차:"
    )
    chain = prompt | llm | StrOutputParser()
    return chain.stream({"topic": topic, "core_toc": core_toc, "context": full_text[:10000]})

def generate_toc(topic, documents):
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-pro", temperature=0.3, google_api_key=gemini_api_key, max_retries=LLM_SDK_MAX_RETRIES)
    full_text = "\n\n".join([doc.page_content for doc in documents])
    prompt = PromptTemplate.from_template(
        "당신은 전문 제안서 컨설턴트입니다. 아래 '주제'와 '참고 자료'를 바탕으로 전문적인 제안서 목차를 생성해주세요.\n"
        "**[규칙]**\n"
        "1. 목차는 서론, 본론(현황분석, 해결과제, 세부 제안), 결론 등의 구조를 가집니다.\n"
        "2. **[매우 중요]** 목차는 반드시 '1. 제안 개요', '2. 시스템 구성 방안'과 같은 최상위 레벨(`X.`)부터 시작해야 합니다. 절대 '1.1.'이나 '1.1.1'과 같은 하위 레벨부터 시작하지 마세요.\n"
        "3. 각 항목은 번호를 붙여주세요. (예: 1. 서론, 2. 현황 분석, ...)\n"
        "4. **[매우 중요]** `-`나 `*` 같은 글머리 기호는 절대 사용하지 마세요.\n"
        "5. **[절대 금지]** 목차 항목에 📌, 🎯 와 같은 이모지(Emoji)나 특수기호는 절대 사용하지 마세요.\n\n"
        "### 주제:\n{topic}\n\n### 참고 자료:\n{context}\n\n### 생성할 목차:"
    )
    chain = prompt | llm | StrOutputParser()
    return chain.stream({"topic": topic, "context": full_text[:10000]})

def generate_section_content_openai(section, topic, toc, documents, page_count, total_sections, attempt_timeout=LLM_TIMEOUT_SECONDS):
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-pro", temperature=0.6, google_api_key=gemini_api_key, timeout=attempt_timeout, max_retries=LLM_SDK_MAX_RETRIES)
    full_text = "\n\n".join([doc.page_content for doc in documents]) if documents else "제공된 참고 자료 없음"
    section_pages = max(page_count / total_sections, 1)
    section_words = int(section_pages * 500)
    min_words = int(section_words * 0.9)
    prompt = PromptTemplate.from_template(
        "당신은 해당 분야의 최고 전문가입니다. 아래 주어진 정보를 바탕으로, 제안서의 '{section}' 섹션에 들어갈 상세 내용을 작성해주세요.\n"
        "**[중요 원칙]**\n"
        "1. 본문 내용은 반드시 상위 목차인 '{section}'에 대한 **논리적 근거와 구체적인 설명**을 포함해야 합니다.\n"
        "2. **[분량 목표]** 이 섹션은 **최소 {min_words} 단어 이상, 약 {section_words} 단어 분량**이 되도록 작성해야 합니다. 이는 전체 제안서가 약 {page_count} 페이지가 되기 위한 필수 조건입니다.\n"
        "3. **[분량 확보 방안]** 만약 내용이 부족하다면, 주장을 뒷받침하는 **구체적인 가상 예시, 통계 자료 인용, 기대 효과, 단계별 실행 방안** 등을 추가하여 반드시 목표 분량을 채워야 합니다. '더 이상 쓸 내용이 없습니다' 와 같이 절대 답변을 중단하지 마세요.\n"
        "4. 다른 섹션의 내용은 절대 작성하지 말고, 오직 '{section}' 부분의 본문만 작성해야 합니다.\n"
        "5. **절대로 목차 제목인 '{section}'을 다시 반복하지 말고, 본문 내용부터 바로 시작해주세요.**\n\n"
        "### 전체 제안서 주제:\n{topic}\n\n"
        "### 전체 목차 (현재 작성 중인 섹션은 '{section}'):\n{toc}\n\n"
        "### 핵심 참고 자료 (업로드된 문서):\n{context}\n\n"
        "### 작성할 '{section}' 섹션의 상세 내용 (최소 {min_words} 단어 이상):"
    )
    chain = prompt | llm | StrOutputParser()
    return chain.invoke({
        "section": section,
        "topic": topic,
        "toc": toc,
        "context": full_text[:8000],
        "page_count": page_count,
        "section_words": section_words,
        "min_words": min_words
    })

def generate_section_with_citations(section, topic, toc, documents, search_tool, status_text, page_count, total_sections, attempt_timeout=LLM_TIMEOUT_SECONDS):
    if not (google_api_key and google_cse_id):
        st.warning("Google API 키 또는 검색엔진 ID가 없어 출처 표기 기능을 사용할 수 없습니다. 일반 모드로 생성합니다.")
        return generate_section_content_openai(section, topic, toc, documents, page_count, total_sections, attempt_timeout=attempt_timeout), ""

    uploaded_context = "\n\n".join([doc.page_content for doc in documents]) if documents else "제공된 참고 자료 없음"
    status_text.text(f"🔄 ({st.session_state.current_section_index + 1}/{st.session_state.total_sections}) '{section}' 근거 자료 검색 중 (법령, 논문, 뉴스)...")
    queries = [
        f"{topic} {section}",
        f"{topic} {section} 법령 OR 규정",
        f"{topic} {section} 연구 OR 논문",
        f"{topic} {section} 공고 OR 보고서"
    ]
    search_results = []
    endpoint = "https://www.googleapis.com/customsearch/v1"
    for query in queries:
        try:
            params = {'key': google_api_key, 'cx': google_cse_id, 'q': query, 'num': 2}
            response = requests.get(endpoint, params=params, timeout=10)
            response.raise_for_status()
            raw_results = response.json().get("items", [])
            for res in raw_results:
                search_results.append({'title': res.get('title'), 'link': res.get('link'), 'snippet': res.get('snippet')})
        except requests.exceptions.RequestException as e:
            st.warning(f"'{query}' 검색 중 네트워크 오류 발생: {e}")
        except Exception as e:
            st.warning(f"'{query}' 검색 결과 처리 중 오류 발생: {e}")
    
    unique_results = {res['link']: res for res in search_results}.values()
    if not unique_results:
        status_text.text(f"🔄 ({st.session_state.current_section_index + 1}/{st.session_state.total_sections}) '{section}' 관련 웹 자료 없음. 내용 생성 중...")
        return generate_section_content_openai(section, topic, toc, documents, page_count, total_sections, attempt_timeout=attempt_timeout), ""
    
    retrieved_context = ""
    for i, res in enumerate(unique_results):
        retrieved_context += f"### 검색자료_{i+1}\n- 제목: {res.get('title', 'N/A')}\n- 링크: {res.get('link', 'N/A')}\n- 요약: {res.get('snippet', 'N/A')}\n\n"

    status_text.text(f"🔄 ({st.session_state.current_section_index + 1}/{st.session_state.total_sections}) '{section}' 자료 분석 및 내용 생성 중...")
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-pro", temperature=0.6, google_api_key=gemini_api_key, timeout=attempt_timeout, max_retries=LLM_SDK_MAX_RETRIES)
    prompt = PromptTemplate.from_template(
        "당신은 최고 수준의 전문 리서처이자 제안서 작성 전문가입니다. 당신의 임무는 '{section}'에 대한 논점을 제시하고, 검색된 자료를 바탕으로 그 논점을 뒷받침하는 근거를 제시하는 것입니다.\n\n"
        "**[작업 절차 및 규칙]**\n"
        "1.  **분량 준수:** 전체 제안서는 약 {page_count} 페이지이며, 총 {total_sections}개의 섹션으로 구성됩니다. 이 섹션은 약 {section_pages} 페이지 분량으로 작성해주세요.\n"
        "2.  **논점 제시:** 먼저 '{section}'의 핵심 주장을 한두 문장으로 제시합니다.\n"
        "3.  **근거 분석:** 제공된 **[검색 자료]** 각각의 제목과 요약을 분석하여, 자료의 유형(예: 일반 웹 기사, 뉴스, 법령, 정부 공고, 연구 논문 등)을 추론합니다.\n"
        "4.  **본문 작성 및 인용:** 분석한 자료를 근거로 사용하여 논점을 뒷받침하는 상세한 본문을 작성합니다. 자료를 인용한 문장 끝에는 `[출처:n]` 형식으로 번호를 반드시 표기합니다.\n"
        "5.  **참고 자료 목록화:** 본문 작성이 끝나면, `---` 구분선 아래에 **'참고 자료'**라는 제목으로 인용한 자료의 목록을 정리합니다. 이때, **추론한 자료 유형에 맞는 형식**을 아래 예시처럼 적용해야 합니다.\n\n"
        "**[출처 표기 형식 예시]**\n"
        "- **일반 웹/뉴스:** `n. [출처] **{{제목}}**, {{링크}}`\n"
        "- **법령/공고:** `n. [출처] **{{제목 \"핵심 내용 요약\"}}**, {{문서 번호 또는 날짜}}, {{링크}}` (문서 번호나 날짜가 제목/요약에 있다면 포함)\n"
        "- **논문/보고서:** `n. [출처] **{{제목}}**, {{저자 또는 기관}}, {{링크}}` (저자나 기관 정보가 있다면 포함)\n\n"
        "**[중요]**\n"
        "- **[업로드된 참고 자료]**는 당신의 배경 지식입니다. 본문 작성 시 참고하되, 직접적인 `[출처]` 표기는 하지 마세요.\n"
        "- 다른 섹션의 내용은 절대 작성하지 말고, 오직 '{section}' 부분의 본문과 참고 자료 목록만 작성해야 합니다.\n\n"
        "---\n**[업로드된 참고 자료 (배경지식용)]**\n{uploaded_context}\n\n"
        "---\n**[검색 자료 (분석 및 인용 대상)]**\n{retrieved_context}\n\n"
        "---\n**[전체 제안서의 맥락]**\n- 전체 주제: {topic}\n- 전체 목차: {toc}\n\n"
        "---\n**['{section}' 섹션 내용 작성 시작]**"
    )
    chain = prompt | llm | StrOutputParser()
    response_text = chain.invoke({
        "section": section, "uploaded_context": uploaded_context[:4000], "retrieved_context": retrieved_context, 
        "topic": topic, "toc": toc, "page_count": page_count, "total_sections": total_sections, 
        "section_pages": round(page_count / total_sections, 1)
    })
    
    if "---" in response_text:
        parts = response_text.split("---", 1)
        content, citations = parts[0].strip(), parts[1].strip()
    else:
        content, citations = response_text.strip(), ""
    return content, citations

# ==============================================================================
# 4단계를 위한 백엔드 함수들
# ==============================================================================

def clamp_progress(value):
    return max(0.0, min(1.0, value))

def update_live_progress(status_text, progress_bar, message, progress_value):
    safe_progress = clamp_progress(progress_value)
    percent = int(round(safe_progress * 100))
    status_text.text(f"🔄 {message} ({percent}%/100%)")
    progress_bar.progress(safe_progress)

def append_enhancement_log(log_container, logs, message):
    timestamp = datetime.datetime.now().strftime("%H:%M:%S")
    logs.append(f"[{timestamp}] {message}")
    visible_logs = logs[-ENHANCEMENT_LOG_LIMIT:]
    log_container.text("실행 로그\n" + "\n".join(visible_logs))

def execute_with_timeout(task, *args, timeout_seconds=LLM_TIMEOUT_SECONDS, on_tick=None, **kwargs):
    """
    task를 별도 스레드에서 실행하고 timeout_seconds 내에 결과를 수립한다.
    task에 attempt_timeout이라는 keyword를 주입해 LLM SDK의 자체 timeout도 같이 걸도록 한다.
    """
    if "attempt_timeout" not in kwargs:
        kwargs["attempt_timeout"] = timeout_seconds
    executor = ThreadPoolExecutor(max_workers=1)
    future = executor.submit(task, *args, **kwargs)
    start_time = time.monotonic()

    try:
        while True:
            try:
                return future.result(timeout=1)
            except FutureTimeoutError:
                elapsed = time.monotonic() - start_time
                if on_tick:
                    on_tick(elapsed, timeout_seconds)
                if elapsed >= timeout_seconds:
                    future.cancel()
                    raise TimeoutError(f"{timeout_seconds}초 제한 시간 초과")
    finally:
        executor.shutdown(wait=False, cancel_futures=True)

def execute_with_retries(task, *args, timeouts=None, on_attempt=None, on_wait=None, **kwargs):
    timeout_sequence = list(timeouts) if timeouts else list(LLM_RETRY_TIMEOUTS)
    if not timeout_sequence:
        timeout_sequence = [LLM_TIMEOUT_SECONDS]
    total_attempts = len(timeout_sequence)
    last_error = None

    for attempt_index, attempt_timeout in enumerate(timeout_sequence, start=1):
        if on_attempt:
            on_attempt(attempt_index, total_attempts, attempt_timeout)
        try:
            tick = None
            if on_wait:
                def tick(elapsed, timeout, current_attempt=attempt_index, total=total_attempts):
                    on_wait(elapsed, timeout, current_attempt, total)
            return execute_with_timeout(task, *args, timeout_seconds=attempt_timeout, on_tick=tick, **kwargs)
        except TimeoutError as exc:
            last_error = exc
            _notify_llm_retry("timeout", attempt_index, total_attempts, exc)
            continue
        except Exception as exc:
            # SDK 내부 에러(네트워크·쿠타·JSON 오류 등)도 동일하게 재시도
            last_error = exc
            # Gemini 504/DeadlineExceeded는 사용자에게 즉시 노출 (stderr 로그가 UI에 안 보이므로)
            err_text = f"{type(exc).__name__}: {exc}"
            if "DeadlineExceeded" in err_text or "504" in err_text:
                _notify_llm_retry("deadline_504", attempt_index, total_attempts, exc)
            else:
                _notify_llm_retry("error", attempt_index, total_attempts, exc)
            continue

    raise last_error if last_error else TimeoutError("재시도 정책이 비어 있습니다.")


def _notify_llm_retry(kind, attempt_index, total_attempts, exc):
    """LLM 재시도 발생 시 Streamlit UI에 토스트로 노출. Streamlit 컨텍스트 밖에서는 stderr로만 출력."""
    msg_map = {
        "timeout": f"⏱ LLM 응답이 시도 {attempt_index}/{total_attempts}에서 타임아웃되었습니다. 다음 시도로 진행합니다.",
        "deadline_504": f"🚨 Gemini 504(DeadlineExceeded) 시도 {attempt_index}/{total_attempts} — Google 측 응답 지연. 다음 시도로 넘어갑니다.",
        "error": f"⚠️ LLM 오류 시도 {attempt_index}/{total_attempts}: {type(exc).__name__}",
    }
    msg = msg_map.get(kind, str(exc))
    # 콘솔 로깅 (백엔드 운영자에게 즉시 가시화)
    _lg = get_logger("LLM")
    log_level = "error" if kind in ("deadline_504",) else "warning"
    getattr(_lg, log_level)(f"retry={kind} attempt={attempt_index}/{total_attempts} err={type(exc).__name__}: {exc}")
    try:
        st.toast(msg, icon="🚨" if kind == "deadline_504" else "⚠️")
        if kind == "deadline_504" and attempt_index == total_attempts:
            st.warning(
                "Gemini API가 반복적으로 504(DeadlineExceeded)를 반환했습니다. 이 섹션은 폴백(원문 유지)으로 처리됩니다. "
                "잠시 후 '🔁 이어서 재실행' 버튼을 사용하시면 폴백 섹션만 다시 시도합니다."
            )
    except Exception:
        pass

def call_llm_api_for_body_processing(prompt_text, attempt_timeout=LLM_TIMEOUT_SECONDS):
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-pro", temperature=0.5, google_api_key=gemini_api_key, timeout=attempt_timeout, max_retries=LLM_SDK_MAX_RETRIES)
    return llm.invoke(prompt_text).content

def call_llm_api_for_toc_suggestion(prompt_text, attempt_timeout=LLM_TIMEOUT_SECONDS):
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-pro", temperature=0.4, google_api_key=gemini_api_key, timeout=attempt_timeout, max_retries=LLM_SDK_MAX_RETRIES)
    response_str = llm.invoke(prompt_text).content.strip()
    try:
        if response_str.startswith("```json"):
            response_str = response_str.split("```json\n", 1)[1].split("```")[0]
        return json.loads(response_str)
    except (json.JSONDecodeError, Exception) as e:
        print(f"JSON 처리 오류: {e}\n원본 응답: {response_str}")
        return {"suggested_title": response_str}

def clean_text(text):
    if not isinstance(text, str): return ""
    cleaned = re.sub(r'[^\w\s\d.,():%가-힣#\n-"‘’“”]', '', text)
    return cleaned.replace('\ufffd', '')

def remove_markdown_headings(text):
    if not isinstance(text, str): return ""
    return re.sub(r'^\s*#{1,3}\s', '', text, flags=re.MULTILINE)

def parse_proposal_string(proposal_text):
    lines = proposal_text.split('\n')
    toc_structure, body_map = [], {}
    current_heading, content_buffer = None, []
    heading_pattern = re.compile(r'^(#{1,3})\s+([\d\.]*\d+\.)\s*(.*)')

    for line in lines:
        match = heading_pattern.match(line)
        if match:
            if current_heading: body_map[current_heading['original_line']] = "\n".join(content_buffer).strip()
            level, number, text = len(match.group(1)), match.group(2).strip(), match.group(3).strip()
            current_heading = {"text": text, "level": level, "original_line": line, "number": number}
            toc_structure.append(current_heading)
            content_buffer = []
        elif line.strip().startswith('# ') and not heading_pattern.match(line):
            if current_heading: body_map[current_heading['original_line']] = "\n".join(content_buffer).strip()
            current_heading = {"text": line.lstrip('# ').strip(), "level": 1, "original_line": line, "number": ""}
            toc_structure.append(current_heading)
            content_buffer = []
        elif current_heading:
            content_buffer.append(line)
    if current_heading: body_map[current_heading['original_line']] = "\n".join(content_buffer).strip()
    return toc_structure, body_map

def process_body_section(text, review_criteria, page_count, on_wait=None, on_attempt=None):
    if not text or not text.strip():
        return "", "본문이 비어 있어 건너뜁니다."
    prompt_body = f"""
    당신은 최고의 제안서 편집 전문가입니다. 아래 [제안서 섹션 내용]을 [품질 검토 체크리스트]에 따라 분석하고 개선해주세요.
    
    [품질 검토 체크리스트]
    {review_criteria}

    [작업 지시]
    1. **분량 조절:** 전체 제안서는 약 {page_count} 페이지입니다. 이 점을 고려하여 내용의 상세도를 조절하되, 절대 내용을 임의로 삭제하거나 과도하게 요약하지 마세요.
    2. **내용 개선:** 체크리스트의 모든 항목을 충족하도록, 문맥을 해치지 않는 선에서 문장을 다듬고, 논리를 보강하며, 어조를 통일합니다.
    3. **구조 유지:** [매우 중요] 원본의 줄바꿈과 단락 구조를 반드시 그대로 유지해야 합니다. 절대 여러 문단을 하나로 합치지 마세요.
    4. **출력 형식:** 수정이 완료된 **최종 섹션 본문(Full Text)만** 출력합니다. 분석 과정이나 설명, 목차 제목은 절대 포함하지 마세요.

    [제안서 섹션 내용]
    {text}
    
    [개선된 섹션 최종본]
    """
    try:
        improved_text = execute_with_retries(
            call_llm_api_for_body_processing,
            prompt_body,
            timeouts=LLM_RETRY_TIMEOUTS,
            on_attempt=on_attempt,
            on_wait=on_wait,
        )
    except TimeoutError:
        total_budget = sum(LLM_RETRY_TIMEOUTS)
        st.warning(f"⚠️ 본문 개선이 {len(LLM_RETRY_TIMEOUTS)}회 재시도(총 {total_budget}초) 안에 끝나지 않아 원문을 유지합니다.")
        return text, f"본문 개선 재시도 모두 시간 초과로 원문 유지"
    except Exception as e:
        st.warning(f"⚠️ 본문 개선 중 오류가 발생해 원문을 유지합니다: {e}")
        return text, f"본문 개선 실패로 원문 유지: {e}"

    if "sorry" in improved_text.lower() or "죄송합니다" in improved_text:
        st.warning(f"⚠️ AI가 일부 본문 내용 개선을 거부했습니다. 원본을 유지합니다.")
        return text, "AI가 본문 개선을 거부해 원문 유지"

    cleaned_text = clean_text(improved_text)
    if not cleaned_text.strip():
        st.warning("⚠️ 본문 개선 결과가 비어 있어 원문을 유지합니다.")
        return text, "본문 개선 응답이 비어 원문 유지"

    return cleaned_text, None

def suggest_refined_toc_title(item, full_content, on_wait=None, on_attempt=None):
    prompt_toc = f"""
    당신은 제안서의 논리 구조를 완성하는 최고 수준의 편집 전문가입니다. 아래 [본문 내용]을 분석하여, 기존 [목차 제목]을 더 구체적이고 논리적인 제목으로 업그레이드해주세요.

    [작업 절차]
    1. **핵심 개념 추출:** [본문 내용]의 핵심 개념을 1~2개의 키워드로 추출합니다. (예: '정보 관리 혁신', '디지털 전환 전략')
    2. **논리적 재조합:** 추출한 [핵심 개념]과 기존 [목차 제목]의 의미를 모두 포함하여, **문맥상 가장 자연스럽고 논리적인 흐름을 갖는 새로운 제목**을 생성합니다.
    3. **[매우 중요]** 생성된 제목은 **반드시 한 문장으로 간결하게** 만들어야 하며, 절대 줄바꿈을 포함해서는 안 됩니다.
    4. **[절대 금지]** 생성하는 제목에는 이모지(Emoji)나 특수기호를 절대 포함하지 마세요.
    5. **[중요]** 결과는 반드시 아래 JSON 형식으로만 응답해야 하며, 업그레이드된 제목만 포함하세요.

    ---
    [좋은 예시]
    - 기존 제목: "BIM의 역할과 중요성"
    - 핵심 개념: "정보 관리 혁신", "디지털 전환 전략"
    - 좋은 결과: "BIM의 정보 관리 혁신 및 디지털 전환 전략"
    ---

    [기존 목차 제목]
    {item['text']}

    [본문 내용]
    {full_content[:8000]}

    [출력할 JSON]
    ```json
    {{
      "suggested_title": "여기에 간결하고 논리적으로 재조합된 새 목차 제목"
    }}
    ```
    """

    try:
        suggestion = execute_with_retries(
            call_llm_api_for_toc_suggestion,
            prompt_toc,
            timeouts=LLM_RETRY_TIMEOUTS,
            on_attempt=on_attempt,
            on_wait=on_wait,
        )
    except TimeoutError:
        total_budget = sum(LLM_RETRY_TIMEOUTS)
        return item['text'], f"목차 강화 {len(LLM_RETRY_TIMEOUTS)}회 재시도(총 {total_budget}초) 초과로 기존 제목 유지"
    except Exception as e:
        return item['text'], f"목차 강화 실패로 기존 제목 유지: {e}"

    if suggestion and "suggested_title" in suggestion:
        cleaned_title = clean_text(suggestion["suggested_title"]).replace('\n', ' ').strip()
        if cleaned_title:
            return cleaned_title, None

    return item['text'], "목차 강화 응답이 비어 기존 제목 유지"

def refine_table_of_contents(toc_structure, body_map, on_item_start=None, on_item_wait=None, on_item_done=None, on_item_attempt=None, cached_titles=None):
    """cached_titles: {original_line: 강화된 제목} - 이전 진행분 재사용."""
    cached_titles = cached_titles or {}
    refined_toc = []
    numbered_items = [item for item in toc_structure if item.get("number")]
    total_numbered_items = len(numbered_items)
    processed_numbered_items = 0

    for i, item in enumerate(toc_structure):
        if not item.get("number"):
            refined_toc.append(item.copy())
            continue

        processed_numbered_items += 1
        if on_item_start:
            on_item_start(processed_numbered_items, total_numbered_items, item)

        refined_item = item.copy()

        # --- (P1-2) 캐시된 강화 결과가 있으면 LLM 호출 생략 ---
        cached_title = cached_titles.get(item["original_line"])
        if cached_title:
            refined_item["text"] = cached_title
            if on_item_done:
                on_item_done(processed_numbered_items, total_numbered_items, refined_item, "이전 진행분 재사용")
            refined_toc.append(refined_item)
            continue

        child_content_parts = [body_map.get(item["original_line"], "")]
        for next_item in toc_structure[i+1:]:
            if next_item["level"] > item["level"]:
                child_content_parts.append(body_map.get(next_item["original_line"], ""))
            else:
                break

        full_content = "\n\n".join(filter(None, child_content_parts))
        if not full_content.strip():
            if on_item_done:
                on_item_done(processed_numbered_items, total_numbered_items, refined_item, "본문이 없어 기존 제목 유지")
            refined_toc.append(refined_item)
            continue

        suggested_title, warning_message = suggest_refined_toc_title(
            item,
            full_content,
            on_wait=lambda elapsed, timeout, attempt_index, total_attempts, current_item=item, current_index=processed_numbered_items, total_items=total_numbered_items: on_item_wait(current_index, total_items, current_item, elapsed, timeout, attempt_index, total_attempts) if on_item_wait else None,
            on_attempt=lambda attempt_index, total_attempts, attempt_timeout, current_item=item, current_index=processed_numbered_items, total_items=total_numbered_items: on_item_attempt(current_index, total_items, current_item, attempt_index, total_attempts, attempt_timeout) if on_item_attempt else None,
        )
        refined_item["text"] = suggested_title

        if on_item_done:
            on_item_done(processed_numbered_items, total_numbered_items, refined_item, warning_message)

        refined_toc.append(refined_item)

    return refined_toc

def reassemble_proposal_string(toc_structure, processed_sections):
    """
    [최종 수정 기능] 
    1. 문서 최상위 제목은 포함
    2. 목차 리스트는 생략
    3. '1. 사업 개요'부터 본문 시작
    4. 목차 제목에 포함된 불필요한 줄바꿈(↵) 기호 제거
    """
    final_content_parts = []

    # 1. 메인 제목(레벨 1, 번호 없음)을 찾아 가장 먼저 추가합니다.
    if toc_structure and toc_structure[0]['level'] == 1 and not toc_structure[0].get("number"):
        main_title_item = toc_structure[0]
        # 제목 텍스트에서 불필요한 줄바꿈 문자(↵)를 제거합니다.
        clean_title = main_title_item['text'].replace('\n', ' ').strip()
        main_title_heading = f"# {clean_title}"
        final_content_parts.append(main_title_heading)

    # 2. 본문이 시작되는 첫 번째 '번호 있는' 목차 항목의 위치를 찾습니다.
    start_index = -1
    for i, item in enumerate(toc_structure):
        if item.get("number"):
            start_index = i
            break

    # 3. 본문 시작 위치부터 끝까지의 내용을 순서대로 추가합니다.
    if start_index != -1:
        content_toc = toc_structure[start_index:]
        for item in content_toc:
            number = item.get("number", "")
            text = item.get("text", "")
            
            # 각 목차 항목 텍스트에서 불필요한 줄바꿈 문자(↵)를 제거합니다.
            clean_text = text.replace('\n', ' ').strip()
            
            full_heading_text = f"{number} {clean_text}".strip()
            new_heading_line = f"{'#' * item['level']} {full_heading_text}"
            
            final_content_parts.append(new_heading_line)
            
            body_text = processed_sections.get(item["original_line"])
            if body_text:
                final_content_parts.append(f"\n{body_text}")
    else:
        # 번호 있는 본문을 찾지 못한 경우에 대한 예외 처리
        print("경고: 번호가 있는 목차 항목을 찾지 못했습니다.")

    return "\n\n".join(final_content_parts)

def enhance_proposal(draft_text, review_criteria, page_count, status_text, progress_bar, log_container, project_id=None, resume=True):
    logs = []
    try:
        update_live_progress(status_text, progress_bar, "[1/4] 문서 구조 분석 중", 0.03)
        append_enhancement_log(log_container, logs, "최종 검토를 시작했습니다.")
        original_toc, body_map = parse_proposal_string(draft_text)

        total_body_sections = len(body_map)
        numbered_toc_sections = len([item for item in original_toc if item.get("number")])
        append_enhancement_log(log_container, logs, f"분석 완료: 본문 {total_body_sections}개, 목차 {numbered_toc_sections}개 항목을 처리합니다.")

        # --- (P1-2) 이전 진행분 로드 ---
        body_progress = load_progress_map(project_id, "stage4_enhance_body") if resume else {}
        toc_progress = load_progress_map(project_id, "stage4_enhance_toc") if resume else {}
        if body_progress or toc_progress:
            append_enhancement_log(
                log_container, logs,
                f"이전 진행분 발견: 본문 {len(body_progress)}/{total_body_sections}건, 목차 {len(toc_progress)}/{numbered_toc_sections}건 → 이어서 진행합니다.",
            )

        update_live_progress(status_text, progress_bar, "[2/4] 본문 내용 개선 준비 중", 0.05)
        processed_body_map = {}

        if total_body_sections == 0:
            append_enhancement_log(log_container, logs, "개선할 본문이 없어 본문 단계는 건너뜁니다.")

        for i, (original_line, content) in enumerate(body_map.items()):
            title_text = original_line.lstrip('# ').strip()
            body_start = 0.05 + (i / max(total_body_sections, 1)) * 0.65
            body_end = 0.05 + ((i + 1) / max(total_body_sections, 1)) * 0.65

            # --- (P1-2) 이미 처리한 섹션은 DB 결과 재사용 ---
            cached = body_progress.get(original_line)
            if cached and cached.get("payload") is not None:
                processed_body_map[original_line] = cached["payload"]
                update_live_progress(
                    status_text, progress_bar,
                    f"[2/4] 본문 재사용 {i+1}/{total_body_sections}: '{title_text}' (이전 결과)",
                    body_end,
                )
                append_enhancement_log(
                    log_container, logs,
                    f"본문 {i+1}/{total_body_sections} 재사용: '{title_text}' (status={cached.get('status')})",
                )
                continue

            update_live_progress(status_text, progress_bar, f"[2/4] 본문 개선 중 {i+1}/{total_body_sections}: '{title_text}'", body_start)
            append_enhancement_log(log_container, logs, f"본문 {i+1}/{total_body_sections} 시작: '{title_text}'")

            section_start_ts = time.monotonic()
            improved_content, warning_message = process_body_section(
                content,
                review_criteria,
                page_count,
                on_wait=lambda elapsed, timeout, attempt_index, total_attempts, current_title=title_text, start=body_start, end=body_end, current_index=i+1, total_items=total_body_sections:
                    update_live_progress(
                        status_text,
                        progress_bar,
                        f"[2/4] 본문 개선 대기 중 {current_index}/{total_items} (재시도 {attempt_index}/{total_attempts}, {int(timeout)}초 제한): '{current_title}' ({int(min(elapsed, timeout))}초 경과)",
                        start + (min(elapsed / max(timeout, 1), 0.9) * (end - start)),
                    ),
                on_attempt=lambda attempt_index, total_attempts, attempt_timeout, current_title=title_text, current_index=i+1, total_items=total_body_sections:
                    append_enhancement_log(
                        log_container,
                        logs,
                        f"본문 {current_index}/{total_items} '{current_title}' 시도 {attempt_index}/{total_attempts} 시작 ({attempt_timeout}초 제한)",
                    ),
            )
            elapsed_sec = time.monotonic() - section_start_ts
            processed_body_map[original_line] = improved_content

            # --- (P1-2) 섹션 1개 끝나는 즉시 DB에 저장 → 중단되어도 다음 번에 재사용 ---
            section_status = "fallback_original" if warning_message else "ok"
            save_progress_item(
                project_id, "stage4_enhance_body", original_line,
                improved_content, status=section_status, elapsed_sec=elapsed_sec,
            )

            update_live_progress(status_text, progress_bar, f"[2/4] 본문 개선 완료 {i+1}/{total_body_sections}: '{title_text}'", body_end)
            append_enhancement_log(
                log_container,
                logs,
                warning_message or f"본문 {i+1}/{total_body_sections} 완료: '{title_text}' ({elapsed_sec:.1f}초)",
            )

        update_live_progress(status_text, progress_bar, "[3/4] 목차 구조 강화 준비 중", 0.70)
        append_enhancement_log(log_container, logs, "목차 강화 단계를 시작합니다.")

        # --- (P1-2) 목차 강화도 캐시 활용 + 결과 영속화 ---
        def _on_toc_done(current_index, total_items, item, warning_message):
            update_live_progress(
                status_text, progress_bar,
                f"[3/4] 목차 강화 완료 {current_index}/{total_items}: '{item['text']}'",
                0.70 + (current_index / max(total_items, 1)) * 0.25,
            )
            append_enhancement_log(
                log_container, logs,
                warning_message or f"목차 {current_index}/{total_items} 완료: '{item['text']}'",
            )
            try:
                save_progress_item(
                    project_id, "stage4_enhance_toc", item.get("original_line", item['text']),
                    item['text'], status="fallback_original" if warning_message else "ok",
                )
            except Exception as _se:
                # 목차 진행분 저장 실패는 사용자에게도 알린다 (재실행/이어서 진행에 영향)
                notify("warning", f"목차 진행분 저장 실패: {_se}. 재실행 시 이 항목은 다시 처리됩니다.")
                append_enhancement_log(log_container, logs, f"⚠️ 목차 진행분 저장 실패: {_se}")

        refined_toc = refine_table_of_contents(
            original_toc,
            processed_body_map,
            on_item_start=lambda current_index, total_items, item:
                (
                    update_live_progress(
                        status_text,
                        progress_bar,
                        f"[3/4] 목차 강화 중 {current_index}/{total_items}: '{item['text']}'",
                        0.70 + ((current_index - 1) / max(total_items, 1)) * 0.25,
                    ),
                    append_enhancement_log(log_container, logs, f"목차 {current_index}/{total_items} 시작: '{item['text']}'")
                ),
            on_item_wait=lambda current_index, total_items, item, elapsed, timeout, attempt_index, total_attempts:
                update_live_progress(
                    status_text,
                    progress_bar,
                    f"[3/4] 목차 강화 대기 중 {current_index}/{total_items} (재시도 {attempt_index}/{total_attempts}, {int(timeout)}초 제한): '{item['text']}' ({int(min(elapsed, timeout))}초 경과)",
                    0.70 + (((current_index - 1) + min(elapsed / max(timeout, 1), 0.9)) / max(total_items, 1)) * 0.25,
                ),
            on_item_attempt=lambda current_index, total_items, item, attempt_index, total_attempts, attempt_timeout:
                append_enhancement_log(
                    log_container,
                    logs,
                    f"목차 {current_index}/{total_items} '{item['text']}' 시도 {attempt_index}/{total_attempts} 시작 ({attempt_timeout}초 제한)",
                ),
            on_item_done=_on_toc_done,
            cached_titles={k: v.get("payload") for k, v in toc_progress.items() if v.get("payload")},
        )

        update_live_progress(status_text, progress_bar, "[4/4] 최종 문서 생성 중", 0.95)
        append_enhancement_log(log_container, logs, "최종 문서를 조립합니다.")
        final_text = reassemble_proposal_string(refined_toc, processed_body_map)

        update_live_progress(status_text, progress_bar, "[4/4] 저장 및 마무리 중", 0.99)
        append_enhancement_log(log_container, logs, "최종 문서 생성이 완료되었습니다.")
        progress_bar.progress(1.0)
        status_text.success("✅ 모든 개선 작업 완료!")
        return final_text
    except Exception as e:
        import traceback
        st.error(f"오류: 제안서 개선 작업 중 문제가 발생했습니다. - {e}")
        append_enhancement_log(log_container, logs, f"❌ 예외 발생: {e}. 지금까지의 진행분은 DB에 저장되어 재진입 시 이어서 처리됩니다.")
        traceback.print_exc()
        return draft_text

def generate_overview_section(full_proposal_text, topic, attempt_timeout=LLM_TIMEOUT_SECONDS):
    llm = ChatGoogleGenerativeAI(model="gemini-2.5-pro", temperature=0.5, google_api_key=gemini_api_key, timeout=attempt_timeout, max_retries=LLM_SDK_MAX_RETRIES)
    text_to_summarize = full_proposal_text.replace("[사업 개요는 최종본 생성 시 자동으로 채워집니다.]", "").strip()
    if "## 목차" in text_to_summarize:
        text_to_summarize = text_to_summarize.split("## 목차", 1)[1]
        if '\n\n#' in text_to_summarize:
            text_to_summarize = text_to_summarize.split('\n\n#', 1)[1]
            text_to_summarize = '#' + text_to_summarize
    prompt = PromptTemplate.from_template(
        "당신은 최고 수준의 제안서 작성 전문가입니다. 아래 제공된 [전체 제안서 본문]을 분석하여, 서론에 해당하는 '사업 개요' 부분을 약 300자 내외의 완결된 문단으로 작성해주세요.\n\n"
        "**[작성 지침]**\n"
        "1. **핵심 요약:** 제안서의 목적, 배경, 핵심 제안 내용을 명확하고 간결하게 요약해야 합니다.\n"
        "2. **전문적 어조:** '-합니다', '-입니다'와 같은 전문적이고 신뢰감 있는 어조를 사용하세요.\n"
        "3. **분량:** 약 300자 내외로 작성하여 서론으로서의 역할을 충실히 수행해야 합니다.\n"
        "4. **출력 형식:** 다른 설명 없이 '사업 개요' 본문 내용만 바로 작성해주세요.\n\n"
        "### 전체 제안서 주제:\n{topic}\n\n"
        "### 전체 제안서 본문 (요약 대상):\n{context}\n\n"
        "### 작성할 '사업 개요' 본문 (약 300자):"
    )
    chain = prompt | llm | StrOutputParser()
    return chain.invoke({"topic": topic, "context": text_to_summarize[:12000]})

def create_docx(content):
    document, is_first_line = Document(), True
    for line in content.split('\n'):
        if line.startswith('# '):
            text = line.lstrip('# ').strip()
            if is_first_line:
                document.add_heading(text, level=0)
                is_first_line = False
            else: document.add_heading(text, level=1)
        elif line.startswith('## '): document.add_heading(line.lstrip('## ').strip(), level=2)
        elif line.startswith('### '): document.add_heading(line.lstrip('### ').strip(), level=3)
        elif line.strip(): document.add_paragraph(line)
        if is_first_line and line.strip().startswith('#'): is_first_line = False
    bio = io.BytesIO()
    document.save(bio)
    return bio.getvalue()

def generate_ppt_slides(text, pages):
    llm, llm_name = get_ppt_llm(gemini_api_key=gemini_api_key, max_retries=LLM_SDK_MAX_RETRIES, temperature=0.5)
    log.info(f"PPT 슬라이드 생성 LLM: {llm_name}")
    try:
        st.caption(f"🧠 PPT 생성 모델: **{llm_name}**")
    except Exception:
        pass
    map_template = """다음은 전체 제안서의 일부 내용입니다. 이 내용을 PPT 슬라이드의 한 부분으로 만들기 위해 핵심 아이디어를 몇 개의 불렛 포인트로 요약해주세요. 내용: {text} 핵심 요약:"""
    map_prompt = PromptTemplate.from_template(map_template)
    reduce_template = """당신은 프레젠테이션 설계 전문가입니다. 아래는 각 파트별로 요약된 제안서의 핵심 내용들입니다. 이 요약본들을 종합하여, 전체적인 논리적 흐름에 맞는 {pages}장 분량의 PPT 슬라이드를 만들어야 합니다. 아래 규칙에 따라 JSON 형식으로만 출력해주세요.
**규칙:**
1. 전체 슬라이드는 'title', '목차', '본문(주요 내용)', '결론', 'Q&A' 슬라이드를 포함하여 총 {pages}장으로 구성합니다.
2. 각 슬라이드는 'title'과 'content' 키를 가집니다.
3. 'content'는 글머리 기호가 포함된 문자열 리스트(['- 항목1', '- 항목2'])로 구성합니다.
4. 'content'의 각 항목은 간결하고 명확해야 합니다.
5. JSON 코드 블록을 사용하지 말고, 순수한 JSON 텍스트만 출력하세요.
**파트별 핵심 내용 요약본:** {text}
**출력할 JSON:**"""
    reduce_prompt = PromptTemplate.from_template(reduce_template)
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=8000, chunk_overlap=400)
    docs = text_splitter.create_documents([text])
    chain = load_summarize_chain(llm, chain_type="map_reduce", map_prompt=map_prompt, combine_prompt=reduce_prompt, verbose=False)
    result_dict = chain.invoke({"input_documents": docs, "pages": pages}, return_only_outputs=True)
    response = result_dict.get('output_text', '')
    if "```json" in response:
        response = response.split("```json\n")[1].split("```")[0]
    return response

def create_ppt_presentation(slides_json, theme_name):
    THEMES = {
        "심플 (파랑)": {"title_color": RGBColor(0, 82, 165), "content_color": RGBColor(89, 89, 89)},
        "비즈니스 (회색)": {"title_color": RGBColor(68, 84, 106), "content_color": RGBColor(102, 102, 102)},
        "크리에이티브 (보라)": {"title_color": RGBColor(102, 0, 255), "content_color": RGBColor(70, 70, 70)},
        "자연 (초록)": {"title_color": RGBColor(0, 128, 0), "content_color": RGBColor(95, 114, 95)},
        "따뜻함 (주황)": {"title_color": RGBColor(255, 102, 0), "content_color": RGBColor(128, 83, 0)},
    }
    selected_theme = THEMES.get(theme_name, THEMES["심플 (파랑)"])
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slides_list = []
    try:
        parsed_data = json.loads(slides_json)
        if isinstance(parsed_data, dict) and "slides" in parsed_data and isinstance(parsed_data["slides"], list):
            slides_list = parsed_data["slides"]
        elif isinstance(parsed_data, list):
            slides_list = parsed_data
        else:
            st.error("AI가 생성한 데이터가 올바른 슬라이드 목록 형식이 아닙니다.")
            st.code(parsed_data)
            return None
    except json.JSONDecodeError:
        st.error("AI가 유효하지 않은 JSON 형식의 데이터를 생성했습니다.")
        st.code(slides_json)
        return None
    for slide_data in slides_list:
        if not isinstance(slide_data, dict):
            continue
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                shape.text_frame.text = slide_data.get("title", "제목 없음")
                shape.text_frame.paragraphs[0].font.color.rgb = selected_theme["title_color"]
                shape.text_frame.paragraphs[0].font.bold = True
            elif shape.placeholder_format.idx == 1:
                tf = shape.text_frame
                tf.clear()
                content_items = slide_data.get("content", [])
                if isinstance(content_items, list):
                    for item in content_items:
                        p = tf.add_paragraph()
                        p.text = item
                        p.level = 0
                        p.font.color.rgb = selected_theme["content_color"]
                else:
                    p = tf.add_paragraph()
                    p.text = str(content_items)
                    p.level = 0
                    p.font.color.rgb = selected_theme["content_color"]
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()

# --- 4. 스트림릿 UI 구성 ---
def reset_all_state():
    keys_to_clear = [key for key in st.session_state if key not in ['Google Search', 'active_tab', 'project_id']]
    for key in keys_to_clear: del st.session_state[key]

init_db()
if not getattr(setup_logging, "_db_logged", False):
    get_logger("DB").info(f"DB 초기화 완료 (WAL): {DB_FILE}")
    setup_logging._db_logged = True

# --- 인증 가드 (모든 페이지 진입 전에 호출) ---
auth_user = require_login()
if st.session_state.get("_logged_user_email") != auth_user["email"]:
    get_logger("AUTH").info(f"로그인 사용자: {auth_user['email']} (role={auth_user['role']})")
    st.session_state._logged_user_email = auth_user["email"]
render_sidebar_user_panel()

# --- 인증 통과 후 메인 헤더 ---
page_header(
    title="🤖 제안서 & 추진계획서 자동 생성 Agent",
    subtitle="5단계 마법사로 제안서·추진계획서를 작성하고, 자동으로 PPT까지 전환합니다.",
    meta="v5.1 · 2026",
)

# --- 외부 API 상태 바 (Google / Gemini / Claude) ---
render_api_status_bar()

if 'selected_project_id' in st.session_state and 'project_loaded' not in st.session_state:
    load_project_into_session(st.session_state.selected_project_id)
    st.session_state.project_loaded = True

if 'Google Search' not in st.session_state: st.session_state['Google Search'] = get_search_tool()
if 'active_tab' not in st.session_state: st.session_state.active_tab = "1단계: 제안서 시작"
if 'generation_stopped' not in st.session_state: st.session_state.generation_stopped = False

# (P1-3) 부팅 시 enhancing 플래그가 남아 있으면 강제로 내린다.
# Streamlit 리로드/세션 재진입 시 LLM 자동 재호출 폭주를 막는다.
if st.session_state.get('enhancing'):
    st.session_state.enhancing = False
    st.warning("이전 강화 작업이 비정상 종료된 것으로 보여 자동 실행을 중지했습니다. '🔁 재실행' 버튼으로 다시 시작하면 이전 진행분을 재사용합니다.")
if st.session_state.get('is_generating'):
    st.session_state.is_generating = False

# (글로벌 CSS는 ui_theme.inject_global_css()에서 이미 주입됨)

tab_names = ["1단계: 제안서 시작", "2단계: 주제/목차 확정", "3단계: 제안서 생성", "4단계: 최종 품질 검증", "5단계: PPT 전환"]
stepper_labels = ["제안서 시작", "주제·목차 확정", "제안서 생성", "최종 품질 검증", "PPT 전환"]
try: active_tab_index = tab_names.index(st.session_state.active_tab)
except ValueError: active_tab_index = 0; st.session_state.active_tab = tab_names[0]

# --- 단계 진행 상태 스테퍼 (현재 세션 상태로 완료 단계 추정) ---
_completed = set()
if st.session_state.get('docs'): _completed.add(0)
if st.session_state.get('finalized_toc'): _completed.add(1)
if st.session_state.get('draft_proposal'): _completed.add(2)
if st.session_state.get('final_proposal'): _completed.add(3)
render_stepper(stepper_labels, current_index=active_tab_index, completed=_completed)

st.session_state.active_tab = st.radio("Navigation", tab_names, index=active_tab_index, horizontal=True, label_visibility="collapsed")

if st.session_state.active_tab == "1단계: 제안서 시작":
    st.header("시작 옵션 선택")
    
    if st.button("🚀 새 제안서 프로젝트 시작하기", use_container_width=True, type="primary"):
        project_id = create_new_project()
        reset_all_state()
        st.session_state.project_id = project_id
        st.success(f"새 프로젝트 (ID: {st.session_state.project_id})가 시작되었습니다. 아래에서 참고 자료를 업로드해주세요.")
        st.rerun()

    if 'project_id' in st.session_state:
        st.info(f"현재 진행 중인 프로젝트 ID: **{st.session_state.project_id}**")
        with st.expander("옵션 1: 새로운 제안서 생성하기 (참고 자료 기반)", expanded=True):
            st.info("제안서 내용의 기반이 될 참고 문서(PDF, TXT)와 핵심 목차(선택)를 업로드하여 새로운 제안서를 만듭니다.")
            core_toc_file = st.file_uploader("핵심 목차(.txt) 파일 (선택)", type="txt", key="new_core_toc")
            ref_files = st.file_uploader("참고 문서 (필수)", type=["pdf", "txt"], accept_multiple_files=True, key="new_ref_files")
            if st.button("학습 및 등록 시작하기", use_container_width=True):
                if not ref_files:
                    st.error("참고 문서는 반드시 1개 이상 업로드해야 합니다."); st.stop()
                if core_toc_file:
                    st.session_state.user_core_toc = core_toc_file.getvalue().decode("utf-8")
                with st.spinner("참고 문서를 분석하고 있습니다..."):
                    texts = []
                    for file in ref_files:
                        try:
                            if file.type == "text/plain": texts.append(io.StringIO(file.getvalue().decode("utf-8")).read())
                            else:
                                pdf_reader = PdfReader(file)
                                texts.append("".join(page.extract_text() for page in pdf_reader.pages if page.extract_text()))
                        except Exception as e:
                            st.error(f"'{file.name}' 파일 처리 중 오류 발생: {e}")
                    docs = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200).create_documents(texts)
                    st.session_state.docs = docs
                    save_stage_result(st.session_state.project_id, "1단계: 자료 업로드", json.dumps([doc.page_content for doc in docs]))

                st.success("학습/등록 완료! 2단계로 이동하여 주제와 목차를 확정하세요.")
                st.session_state.active_tab = "2단계: 주제/목차 확정"
                st.rerun()
    else:
        st.warning("새 제안서 프로젝트를 시작해주세요.")

elif st.session_state.active_tab == "2단계: 주제/목차 확정":
    if 'project_id' not in st.session_state:
        st.warning("먼저 '1단계'에서 새 프로젝트를 시작해주세요.")
        st.stop()
    st.header("주제 및 목차 상세화, 확정")
    if 'docs' not in st.session_state or not st.session_state.docs:
        st.warning("먼저 '1단계'에서 참고 자료를 등록해주세요.")
    
    if 'finalized_topic' not in st.session_state:
        if 'recommendations' not in st.session_state:
            if st.button("학습 내용 기반 AI 주제 추천받기", use_container_width=True, disabled=(not st.session_state.get('docs'))):
                with st.status("AI가 학습 내용을 분석하여 주제명을 추천하고 있습니다...", expanded=True) as status:
                    status.write("1/2: 업로드된 문서를 작은 단위로 요약 중입니다... (맵 단계)")
                    recommendations = get_title_recommendations(st.session_state.docs)
                    status.write("2/2: 요약된 내용들을 종합하여 최종 주제명을 생성 중입니다... (리듀스 단계)")
                    st.session_state.recommendations = recommendations
                    status.update(label="✅ 주제 추천이 완료되었습니다!", state="complete", expanded=False)
                st.rerun()
        else:
            st.subheader("1. AI 추천 주제 선택")
            for topic in st.session_state.recommendations:
                if st.button(topic, use_container_width=True):
                    st.session_state.selected_topic = topic
                    st.rerun()
            if 'selected_topic' in st.session_state:
                st.subheader("2. 주제 수정 및 확정")
                st.text_input("선택된 주제 (수정 가능)", value=st.session_state.selected_topic, key="topic_editor")
                if st.button("이 주제로 최종 확정", type="primary"):
                    st.session_state.finalized_topic = st.session_state.topic_editor
                    update_project_topic(st.session_state.project_id, st.session_state.finalized_topic)
                    save_stage_result(st.session_state.project_id, "2단계: 주제 확정", st.session_state.finalized_topic)
                    st.rerun()
    
    if 'finalized_topic' in st.session_state:
        st.success(f"주제 확정 완료: **{st.session_state.finalized_topic}**")
        st.markdown("---")
        st.subheader("3. 목차 생성 및 수정")

        if 'editable_toc' not in st.session_state and 'finalized_toc' not in st.session_state:
            if 'user_core_toc' in st.session_state:
                st.info("사용자가 업로드한 핵심 목차 구조를 기반으로 세부 목차를 생성합니다.")
                if st.button("세부 목차 생성하기", use_container_width=True, disabled=(not st.session_state.get('docs'))):
                    with st.chat_message("assistant"):
                        response = st.write_stream(generate_detailed_toc(st.session_state.finalized_topic, st.session_state.user_core_toc, st.session_state.docs))
                    lines = response.split('\n')
                    cleaned_lines = [line.strip() for line in lines if line.strip()]
                    cleaned_toc = "\n".join(cleaned_lines)
                    st.session_state.editable_toc = cleaned_toc
                    st.rerun()
            else:
                st.info("AI가 확정된 주제에 맞춰 전체 목차 초안을 생성합니다.")
                if st.button("목차 초안 생성하기", use_container_width=True, disabled=(not st.session_state.get('docs'))):
                    with st.chat_message("assistant"):
                        response = st.write_stream(generate_toc(st.session_state.finalized_topic, st.session_state.docs))
                    lines = response.split('\n')
                    cleaned_lines = [line.strip() for line in lines if line.strip()]
                    cleaned_toc = "\n".join(cleaned_lines)
                    st.session_state.editable_toc = cleaned_toc
                    st.rerun()
        
        if 'editable_toc' in st.session_state or 'finalized_toc' in st.session_state:
            toc_content_to_edit = st.session_state.get('finalized_toc', st.session_state.get('editable_toc', ''))
            st.markdown("""<style>textarea[data-testid="stTextarea"] > div > textarea {font-size: 1.1rem;}</style>""", unsafe_allow_html=True)
            line_count = toc_content_to_edit.count('\n') + 1
            dynamic_height = line_count * 30 + 30
            st.text_area("생성된 목차 (수정 가능)", value=toc_content_to_edit, height=dynamic_height, key="toc_final_editor")
            
            button_label = "이 목차로 업데이트" if 'finalized_toc' in st.session_state else "이 목차로 최종 확정"
            if st.button(button_label, type="primary", use_container_width=True):
                st.session_state.finalized_toc = st.session_state.toc_final_editor
                save_stage_result(st.session_state.project_id, "2단계: 목차 확정", st.session_state.finalized_toc)
                st.success("목차가 성공적으로 저장/업데이트되었습니다.")
                st.rerun()

        if 'finalized_toc' in st.session_state:
            st.info("준비가 되면 '3단계' 탭으로 이동하여 제안서 생성을 시작하세요.")

elif st.session_state.active_tab == "3단계: 제안서 생성":
    if 'project_id' not in st.session_state: st.warning("먼저 '1단계'에서 새 프로젝트를 시작해주세요."); st.stop()
    st.header("제안서 자동 생성")
    if 'finalized_topic' not in st.session_state or 'finalized_toc' not in st.session_state: st.warning("먼저 '2단계'에서 주제와 목차를 모두 확정해주세요."); st.stop()
    
    page_count = st.number_input("목표 페이지 수", min_value=5, max_value=100, value=st.session_state.get('page_count', 15))
    generation_mode = st.selectbox("생성 모드 선택", ["외부 LLM 전용 (안정적, 추천)", "나만의 LLM (로컬 서버 필요)"])
    use_citations = st.checkbox("웹 검색으로 내용 보강 및 [출처] 표기하기 (생성 속도가 느려질 수 있습니다.)")

    st.markdown("---")
    st.info(f"**확정된 주제:** {st.session_state.finalized_topic}")
    with st.expander("✅ 확정된 최종 목차 보기 (클릭)"): st.code(st.session_state.get('finalized_toc', ''), language='text')
    st.markdown("---")

    if st.session_state.get('is_generating'):
        if st.button("🛑 생성 중단하기 (현재 섹션 종료 후 정지)", use_container_width=True, type="primary"):
            st.session_state.generation_stopped = True
            st.session_state.is_generating = False
            st.warning("정지 요청 접수: 진행 중인 섹션이 끝나면 정지합니다. 이미 생성된 섹션은 DB에 저장되어 다음 실행 시 자동으로 재사용됩니다.")
            st.rerun()
    elif 'draft_proposal' not in st.session_state:
        # --- (P1-4) 3단계 진행분 재사용 카드 ---
        stage3_summary = get_progress_summary(st.session_state.project_id, "stage3_body")
        if stage3_summary:
            st.warning(
                f"이전에 생성하다 중단된 본문이 {stage3_summary['count']}개 남아 있습니다 (마지막 {stage3_summary['latest']}). "
                "다시 시작하면 처리된 섹션은 LLM 재호출 없이 재사용됩니다."
            )
            # 현재 목차 기준 누락/중복 검증 리포트
            try:
                _expected_headings = [
                    line.strip() for line in (st.session_state.get('finalized_toc', '') or '').split('\n')
                    if line.strip() and not line.strip().startswith('#')
                ]
                _rep = verify_progress_integrity(st.session_state.project_id, "stage3_body", _expected_headings)
                render_integrity_report(_rep, title="3단계 진행분 검증")
                if _rep["duplicates"]:
                    notify("error", f"3단계 진행분에 중복 섹션 {len(_rep['duplicates'])}건이 감지됐습니다. '처음부터' 버튼을 권장합니다.")
            except Exception as _verr:
                notify("warning", f"3단계 진행분 검증 중 오류: {_verr}")

            rs1, rs2 = st.columns(2)
            with rs1:
                st.caption("▶ '제안서 생성 시작'을 누르면 자동으로 이어서 진행됩니다.")
            with rs2:
                if st.button("🗑 이전 진행분 비우고 처음부터", use_container_width=True, key="reset_stage3_progress"):
                    clear_progress(st.session_state.project_id, "stage3_body")
                    notify("success", "이전 진행분을 비웠습니다. 처음부터 다시 생성됩니다.")
                    st.rerun()
        if st.button("🚀 제안서 생성 시작하기", use_container_width=True):
            if 'draft_proposal' in st.session_state: del st.session_state['draft_proposal']
            if 'citations' in st.session_state: del st.session_state['citations']
            st.session_state.is_generating = True
            st.session_state.generation_stopped = False
            update_project_page_count(st.session_state.project_id, page_count)
            st.rerun()

    if st.session_state.get('is_generating'):
        final_topic, final_toc = st.session_state.finalized_topic, st.session_state.finalized_toc

        # --- 목차 정규화 및 검증 ---
        normalized_toc, toc_warnings = validate_and_normalize_toc(final_toc)
        if toc_warnings:
            with st.expander(f"⚠️ 목차 자동 정규화 항목 ({len(toc_warnings)}건) - 클릭하여 확인", expanded=False):
                for w in toc_warnings:
                    st.caption(f"• {w}")
        final_toc = normalized_toc

        all_headings = [line.strip() for line in final_toc.split('\n') if line.strip() and not line.strip().startswith('#')]
        st.session_state.total_sections = len(all_headings)

        # --- (P1-4) 3단계 진행분 캐시 로드 ---
        stage3_cache = load_progress_map(st.session_state.project_id, "stage3_body")

        st.markdown("### 제안서 생성 진행률")
        progress_bar, status_text = st.progress(0), st.empty()
        full_proposal, all_citations = f"# {final_topic}\n\n", ""
        if final_toc: full_proposal += f"## 목차\n{final_toc}\n\n"
        llm_type = "OpenAI"

        try:
            get_logger("STAGE3").info(
                f"start project={st.session_state.project_id} sections={st.session_state.total_sections} "
                f"page_count={page_count} citations={use_citations}"
            )
            for i, current_heading in enumerate(all_headings):
                st.session_state.current_section_index = i
                if st.session_state.get('generation_stopped'):
                    status_text.warning("사용자 요청으로 정지합니다. 지금까지 처리된 섹션은 DB에 저장되어 다음 실행 시 재사용됩니다.")
                    break
                progress_bar.progress((i + 1) / st.session_state.total_sections)

                is_leaf_node = True
                current_level_match = re.match(r'^(\d+(\.\d+)*)', current_heading)
                if current_level_match and i + 1 < len(all_headings):
                    next_heading = all_headings[i+1]
                    next_level_match = re.match(r'^(\d+(\.\d+)*)', next_heading)
                    if next_level_match and next_level_match.group(1).startswith(current_level_match.group(1) + '.'):
                        is_leaf_node = False

                number_match = re.match(r'^(\d+(?:\.\d+)*\.?)', current_heading)
                if number_match:
                    heading_level = get_heading_level_from_number(number_match.group(1))
                else:
                    heading_level = 1
                full_proposal += f"{'#' * heading_level} {current_heading}\n\n"

                if i == 0 and ("사업 개요" in current_heading or "사업 목적" in current_heading):
                    status_text.text(f"➡️ ({i+1}/{st.session_state.total_sections}) '{current_heading}'는 4단계에서 자동 생성됩니다.")
                    full_proposal += "[사업 개요는 최종본 생성 시 자동으로 채워집니다.]\n\n"
                    continue

                if is_leaf_node:
                    cached_item = stage3_cache.get(current_heading)
                    if cached_item and cached_item.get("payload"):
                        status_text.text(f"♻️ ({i+1}/{st.session_state.total_sections}) '{current_heading}' 이전 결과 재사용")
                        full_proposal += f"{cached_item['payload']}\n\n"
                        continue

                    status_text.text(f"🔄 ({i+1}/{st.session_state.total_sections}) '{current_heading}' 섹션 본문 생성 중... (재시도 정책 {'+'.join(str(t) for t in LLM_RETRY_TIMEOUTS)}초)")
                    generated_content, citations_for_section = "", ""
                    section_start_ts = time.monotonic()
                    try:
                        if use_citations:
                            llm_type = "OpenAI_with_Search"
                            def _task_with_citations(attempt_timeout=LLM_TIMEOUT_SECONDS, **_kw):
                                return generate_section_with_citations(
                                    current_heading, final_topic, final_toc,
                                    st.session_state.get('docs', []), st.session_state['Google Search'],
                                    status_text, page_count, len(all_headings),
                                    attempt_timeout=attempt_timeout,
                                )
                            generated_content, citations_for_section = execute_with_retries(
                                _task_with_citations, timeouts=LLM_RETRY_TIMEOUTS,
                            )
                        else:
                            def _task_no_citations(attempt_timeout=LLM_TIMEOUT_SECONDS, **_kw):
                                return generate_section_content_openai(
                                    current_heading, final_topic, final_toc,
                                    st.session_state.get('docs', []), page_count, len(all_headings),
                                    attempt_timeout=attempt_timeout,
                                )
                            generated_content = execute_with_retries(
                                _task_no_citations, timeouts=LLM_RETRY_TIMEOUTS,
                            )
                    except TimeoutError:
                        total_budget = sum(LLM_RETRY_TIMEOUTS)
                        notify("warning", f"⚠️ '{current_heading}' 생성이 {len(LLM_RETRY_TIMEOUTS)}회 재시도(총 {total_budget}초) 안에 끝나지 않아 빈 본문으로 진행합니다.")
                        generated_content = f"[자동 생성 실패: 시간 초과 - '{current_heading}']"
                    except Exception as e:
                        notify("error", f"'{current_heading}' 섹션 생성 중 오류 발생: {e}")
                        generated_content = f"[오류: '{current_heading}' 생성 실패 - {e}]"

                    cleaned_section_title = re.sub(r"^\d+(\.\d+)*\s*", "", current_heading).strip()
                    temp_content = generated_content.strip()
                    if temp_content.startswith(current_heading):
                        generated_content = temp_content[len(current_heading):].strip()
                    elif temp_content.startswith(cleaned_section_title):
                        generated_content = temp_content[len(cleaned_section_title):].strip()

                    full_proposal += f"{generated_content}\n\n"
                    if citations_for_section: all_citations += f"### {current_heading}\n{citations_for_section}\n\n"

                    # --- (P1-4) 섹션 1개 완료 즉시 DB에 저장 ---
                    elapsed_sec = time.monotonic() - section_start_ts
                    save_progress_item(
                        st.session_state.project_id, "stage3_body", current_heading,
                        generated_content,
                        status="ok" if not generated_content.startswith("[") else "fallback_original",
                        elapsed_sec=elapsed_sec,
                    )
                    get_logger("STAGE3").info(
                        f"section {i+1}/{st.session_state.total_sections} '{current_heading}' "
                        f"saved status={'ok' if not generated_content.startswith('[') else 'fallback'} "
                        f"elapsed={elapsed_sec:.1f}s"
                    )

                else:
                    status_text.text(f"➡️ ({i+1}/{st.session_state.total_sections}) '{current_heading}' 상위 목차 추가 (본문 생성 건너뜀)")
                    full_proposal += "\n"

            if not st.session_state.get('generation_stopped'):
                st.session_state.draft_proposal = full_proposal.strip()
                save_stage_result(st.session_state.project_id, "3단계: 본문 생성", st.session_state.draft_proposal, llm_type)
                if all_citations:
                    st.session_state.citations = all_citations.strip()
                    save_stage_result(st.session_state.project_id, "3단계: 출처 목록", st.session_state.citations)
                # --- 완료 직전 누락/중복 검증 (사용자 요구사항) ---
                completion_report = verify_progress_integrity(
                    st.session_state.project_id, "stage3_body", all_headings
                )
                if completion_report["missing"] or completion_report["duplicates"]:
                    notify("error",
                           f"🚨 3단계 완료 후 검증 실패 — 누락 {len(completion_report['missing'])}건 / "
                           f"중복 {len(completion_report['duplicates'])}건. 아래 리포트를 확인하세요.")
                    render_integrity_report(completion_report, title="3단계 완료 후 검증")
                elif completion_report["fallback"]:
                    notify("warning",
                           f"⚠️ 3단계 완료. 폴백(원문 유지) 섹션 {len(completion_report['fallback'])}건이 있습니다. 재실행으로 재생성을 권장합니다.")
                else:
                    notify("success", f"✅ 3단계 검증 통과 (총 {completion_report['cached']}개 섹션, 누락·중복 0건).")
                # 정상 완료 시 3단계 캐시 정리 (재실행 충돌 방지)
                clear_progress(st.session_state.project_id, "stage3_body")
                st.success("🎉 제안서 초안 생성이 완료되었습니다!")
            else:
                # 부분 본문도 보존: 중단 시점까지의 draft를 저장해 두면 4단계 진입 시 즉시 활용 가능
                if full_proposal.strip():
                    save_stage_result(st.session_state.project_id, "3단계: 본문 생성", full_proposal.strip(), llm_type)
                    notify("info", "중단 시점까지의 부분 초안을 저장했습니다. 다시 '시작'을 누르면 남은 섹션부터 이어서 생성합니다.")
        except Exception as _stage3_exc:
            get_logger("STAGE3").exception(f"치명적 오류: {_stage3_exc}")
            notify("error", f"🚨 3단계 실행 중 치명적 오류: {_stage3_exc}. 진행분은 DB에 보존되어 재실행 시 자동 재사용됩니다.")
        finally:
            # (P1-3) 어떤 경우에도 is_generating 플래그 해제
            st.session_state.is_generating = False
            get_logger("STAGE3").info("end (is_generating=False)")
        st.rerun()

    if 'draft_proposal' in st.session_state:
        st.markdown("---")
        st.subheader("생성된 제안서 초안")
        display_text = remove_markdown_headings(st.session_state.draft_proposal)
        st.text_area("초안 내용", value=display_text, height=400)
        cols = st.columns([1, 1, 1])
        draft_docx = create_docx(st.session_state.draft_proposal)
        cols[0].download_button(label="📥 초안 DOCX 다운로드", data=draft_docx, file_name=f"[초안] {st.session_state.finalized_topic}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
        cols[1].download_button(label="📥 초안 TXT 다운로드", data=display_text.encode('utf-8'), file_name=f"[초안] {st.session_state.finalized_topic}.txt", mime="text/plain", use_container_width=True)
        if st.session_state.get('citations'):
            cols[2].download_button(label="📥 출처 목록 TXT 다운로드", data=st.session_state.citations.encode('utf-8'), file_name=f"[출처] {st.session_state.finalized_topic}.txt", mime="text/plain", use_container_width=True)

        # --- 3단계 재실행 컨트롤 ---
        st.markdown("---")
        st.markdown("#### 🔁 3단계 재실행")
        rr1, rr2 = st.columns(2)
        with rr1:
            if st.button("🔁 이어서 재실행 (누락된 섹션만 보강)", use_container_width=True, key="stage3_rerun_resume"):
                # 누락된 섹션을 기준으로 재실행: draft_proposal과 citations 캐시는 지우되 stage_progress는 보존
                del st.session_state['draft_proposal']
                if 'citations' in st.session_state: del st.session_state['citations']
                st.session_state.is_generating = True
                st.session_state.generation_stopped = False
                notify("info", "이어서 재실행을 시작합니다. 처리된 섹션은 재사용되고, 누락분만 새로 생성됩니다.")
                st.rerun()
        with rr2:
            if st.button("🔄 처음부터 재실행 (캐시 전부 삭제)", use_container_width=True, key="stage3_rerun_full", type="secondary"):
                clear_progress(st.session_state.project_id, "stage3_body")
                del st.session_state['draft_proposal']
                if 'citations' in st.session_state: del st.session_state['citations']
                st.session_state.is_generating = True
                st.session_state.generation_stopped = False
                notify("warning", "캐시를 모두 삭제하고 처음부터 재실행합니다.")
                st.rerun()

elif st.session_state.active_tab == "4단계: 최종 품질 검증":
    st.header("🧐 제안서 최종 품질 검증")
    if 'project_id' in st.session_state:
        project_id = st.session_state.project_id
        project_data, stages = get_project_data(project_id)
        draft_proposal = stages.get('3단계: 본문 생성')
        if project_data and draft_proposal:
            st.info(f"**검토 대상:** {project_data['topic']} (ID: {project_id})")
            # --- (P1-2) 이전 진행분 재사용 카드 ---
            body_progress_summary = get_progress_summary(project_id, "stage4_enhance_body")
            toc_progress_summary = get_progress_summary(project_id, "stage4_enhance_toc")
            if (body_progress_summary or toc_progress_summary) and not st.session_state.get('enhancing'):
                with st.container():
                    st.warning(
                        "이전에 진행하다 중단된 강화 작업이 남아 있습니다. 이어서 진행하면 처리된 섹션은 LLM 재호출 없이 그대로 사용합니다.\n"
                        f"- 본문: {body_progress_summary['count'] if body_progress_summary else 0}건 (마지막 {body_progress_summary['latest'] if body_progress_summary else '-'})\n"
                        f"- 목차: {toc_progress_summary['count'] if toc_progress_summary else 0}건 (마지막 {toc_progress_summary['latest'] if toc_progress_summary else '-'})"
                    )
                    rc1, rc2 = st.columns(2)
                    with rc1:
                        if st.button("▶ 이어서 진행 (이전 결과 재사용)", use_container_width=True, key="resume_stage4"):
                            st.session_state.resume_stage4 = True
                            st.success("다음 강화 실행 시 이전 진행분을 자동으로 이어서 사용합니다.")
                    with rc2:
                        if st.button("🗑 진행분 비우고 처음부터", use_container_width=True, key="reset_stage4_progress"):
                            clear_progress(project_id, "stage4_enhance_body")
                            clear_progress(project_id, "stage4_enhance_toc")
                            st.session_state.resume_stage4 = False
                            st.success("이전 진행분을 모두 비웠습니다.")
                            st.rerun()
                    st.markdown("---")

            if 'draft_loaded_for_review' not in st.session_state:
                st.session_state.draft_loaded_for_review = False

            def _prepare_draft_from_source(source_text, force_overview=False):
                """초안 텍스트를 검토 영역에 올린다. 사업 개요 토큰이 남아 있을 때만 LLM을 호출한다."""
                placeholder = "[사업 개요는 최종본 생성 시 자동으로 채워집니다.]"
                if force_overview or placeholder in source_text:
                    with st.spinner("AI가 초안을 분석하여 '사업 개요'를 생성하고 있습니다..."):
                        overview_content = generate_overview_section(source_text, project_data.get('topic', ''))
                    source_text = source_text.replace(placeholder, overview_content.strip())
                st.session_state.editable_draft_content = source_text
                st.session_state.draft_loaded_for_review = True

            # 최종본이 이미 있으면 재검증 진입로를 항상 노출
            if 'final_proposal' in st.session_state and not st.session_state.draft_loaded_for_review:
                st.subheader("🔁 재검증 시작")
                st.caption("이전 최종본이 저장되어 있습니다. 어떤 기준으로 다시 검증할지 선택하세요.")
                reverify_col1, reverify_col2, reverify_col3 = st.columns(3)
                with reverify_col1:
                    if st.button("📝 마지막 초안으로 다시 검증", use_container_width=True):
                        base = st.session_state.get('editable_draft_content') or draft_proposal
                        _prepare_draft_from_source(base)
                        st.rerun()
                with reverify_col2:
                    if st.button("✨ 최종본을 새 초안으로 다시 검증", use_container_width=True):
                        _prepare_draft_from_source(st.session_state.final_proposal)
                        st.rerun()
                with reverify_col3:
                    if st.button("↩️ 원본 초안으로 되돌려 다시 검증", use_container_width=True):
                        st.session_state.pop('editable_draft_content', None)
                        _prepare_draft_from_source(draft_proposal, force_overview=True)
                        st.rerun()
                st.markdown("---")

            if not st.session_state.draft_loaded_for_review:
                if st.button("📝 초안 불러와서 검토 및 수정하기", use_container_width=True):
                    base = st.session_state.get('editable_draft_content') or draft_proposal
                    _prepare_draft_from_source(base)
                    st.rerun()

            if st.session_state.draft_loaded_for_review:
                st.subheader("1. 초안 검토 및 수정")
                st.text_area("제안서 초안 (직접 수정 가능)", value=st.session_state.editable_draft_content, height=400, key="editable_draft_content_area")
                action_col1, action_col2 = st.columns(2)
                with action_col1:
                    if st.button("💾 현재 수정 내용 저장(검증 없이)", use_container_width=True):
                        st.session_state.editable_draft_content = st.session_state.editable_draft_content_area
                        st.success("수정한 초안이 세션에 저장되었습니다.")
                with action_col2:
                    if st.button("🧹 검토 영역 닫기", use_container_width=True):
                        st.session_state.editable_draft_content = st.session_state.editable_draft_content_area
                        st.session_state.draft_loaded_for_review = False
                        st.rerun()
                st.markdown("---")
                st.subheader("2. AI 최종 검토 실행")
                st.info("위 초안을 바탕으로 AI가 전체적인 품질을 검토하고 개선합니다. (목차, 본문, 어조 등)")

                review_criteria = st.text_area("품질 검토 기준 (AI에게 전달됩니다)",
                    value="""1.  **어조 및 일관성:** 전체적으로 '-합니다', '-입니다'의 일관된 경어체를 사용하는가?
                             2.  **목차와 본문의 일관성:** 각 본문이 상위 목차의 내용을 충실히 설명하고 뒷받침하는가?
                             3.  **상위 목차의 대표성:** 상위 목차가 하위 본문의 핵심 내용을 잘 요약하고 포괄하는가? (AI가 더 구체적인 제목으로 개선)
                             4.  **전문성 및 설득력:** 제안서 전체적으로 전문적인 용어와 어조를 사용하며, 주장에 대한 근거가 명확한가?
                             5.  **가독성:** 불필요한 공백이나 반복적인 표현이 제거되었는가?""",
                    height=250)

                run_label = "🔁 수정된 초안으로 AI 최종 검토 재실행" if 'final_proposal' in st.session_state else "🚀 수정된 초안으로 AI 최종 검토 실행"
                if st.button(run_label, use_container_width=True, type="primary"):
                    st.session_state.editable_draft_content = st.session_state.editable_draft_content_area
                    st.session_state.enhancing = True
                    st.session_state.draft_to_enhance = st.session_state.editable_draft_content_area
                    st.rerun()

            if st.session_state.get('enhancing'):
                progress_bar, status_text = st.progress(0), st.empty()
                st.caption(f"최종 검토 실행 중: AI 호출은 {'+'.join(str(t) for t in LLM_RETRY_TIMEOUTS)}초로 최대 {len(LLM_RETRY_TIMEOUTS)}회 재시도 후에만 원문 유지로 넘어갑니다. 처리된 섹션은 즉시 DB에 저장되어 중단되어도 이어서 진행할 수 있습니다.")
                log_container = st.empty()
                # 진입 시 진행분 알림
                _entry_body = get_progress_summary(project_id, "stage4_enhance_body")
                _entry_toc = get_progress_summary(project_id, "stage4_enhance_toc")
                if _entry_body or _entry_toc:
                    notify("info",
                           f"4단계 진입 — 본문 {(_entry_body or {}).get('count', 0)}건 / "
                           f"목차 {(_entry_toc or {}).get('count', 0)}건 재사용 예정.")
                try:
                    get_logger("STAGE4").info(
                        f"start project={project_id} resume={st.session_state.get('resume_stage4', True)}"
                    )
                    final_proposal_text = enhance_proposal(
                        st.session_state.draft_to_enhance,
                        review_criteria,
                        project_data.get('page_count', 15),
                        status_text,
                        progress_bar,
                        log_container,
                        project_id=project_id,
                        resume=st.session_state.get('resume_stage4', True),
                    )
                    st.session_state.final_proposal = final_proposal_text
                    save_stage_result(st.session_state.project_id, "4단계: 최종본", st.session_state.final_proposal)
                    # 정상 종료 후 폴백/완료 알림
                    _final_body = get_progress_summary(project_id, "stage4_enhance_body") or {}
                    _final_toc = get_progress_summary(project_id, "stage4_enhance_toc") or {}
                    _fallback_total = _final_body.get('fallback_count', 0) + _final_toc.get('fallback_count', 0)
                    if _fallback_total > 0:
                        notify("warning",
                               f"⚠️ 4단계 완료 — 다만 {_fallback_total}개 섹션이 LLM 실패로 원문 유지(폴백)로 처리되었습니다. "
                               "재실행 버튼으로 폴백 섹션만 재시도할 수 있습니다.")
                    else:
                        notify("success", "✅ 4단계 검토 완료. 폴백 없이 모두 성공했습니다.")
                    # 정상 종료 시 진행분은 정리 (재실행 시 불필요한 캐시 방지)
                    clear_progress(project_id, "stage4_enhance_body")
                    clear_progress(project_id, "stage4_enhance_toc")
                    st.session_state.resume_stage4 = False
                except Exception as _stage4_exc:
                    get_logger("STAGE4").exception(f"치명적 오류: {_stage4_exc}")
                    notify("error",
                           f"🚨 4단계 강화 중 치명적 오류: {_stage4_exc}. "
                           "진행분은 DB에 보존되어 재실행 시 자동 재사용됩니다.")
                finally:
                    # 예외/리로드 시에도 플래그를 반드시 해제 → LLM 폭주 방지 (P1-3)
                    st.session_state.enhancing = False
                    if 'draft_to_enhance' in st.session_state: del st.session_state['draft_to_enhance']
                    get_logger("STAGE4").info("end (enhancing=False)")
                st.rerun()

            if 'final_proposal' in st.session_state:
                st.markdown("---")
                st.subheader("개선된 제안서 최종본")
                display_text = remove_markdown_headings(st.session_state.final_proposal)
                st.text_area("최종본 내용", value=display_text, height=400)

                # --- 4단계 재실행 컨트롤 ---
                st.markdown("#### 🔁 4단계 재실행")
                rr1, rr2 = st.columns(2)
                with rr1:
                    if st.button("🔁 이어서 재실행 (캐시 재사용)", use_container_width=True, key="stage4_rerun_resume"):
                        st.session_state.resume_stage4 = True
                        st.session_state.enhancing = True
                        st.session_state.draft_to_enhance = (
                            st.session_state.get('editable_draft_content') or draft_proposal
                        )
                        notify("info", "이어서 재실행합니다. 처리된 섹션은 LLM 재호출 없이 재사용됩니다.")
                        st.rerun()
                with rr2:
                    if st.button("🔄 처음부터 재실행 (캐시 삭제)", use_container_width=True, key="stage4_rerun_full"):
                        clear_progress(project_id, "stage4_enhance_body")
                        clear_progress(project_id, "stage4_enhance_toc")
                        st.session_state.resume_stage4 = False
                        st.session_state.enhancing = True
                        st.session_state.draft_to_enhance = (
                            st.session_state.get('editable_draft_content') or draft_proposal
                        )
                        notify("warning", "캐시를 삭제하고 4단계를 처음부터 다시 실행합니다.")
                        st.rerun()

                reuse_col1, reuse_col2 = st.columns(2)
                with reuse_col1:
                    if st.button("✨ 이 최종본을 새 초안으로 다시 검증", use_container_width=True, key="reuse_final_as_draft"):
                        _prepare_draft_from_source(st.session_state.final_proposal)
                        st.rerun()
                with reuse_col2:
                    if st.button("↩️ 원본 초안으로 되돌려 다시 검증", use_container_width=True, key="reset_to_original_draft"):
                        st.session_state.pop('editable_draft_content', None)
                        _prepare_draft_from_source(draft_proposal, force_overview=True)
                        st.rerun()

                final_docx = create_docx(st.session_state.final_proposal)
                col1, col2 = st.columns(2)
                with col1:
                    st.download_button(label="📥 최종본 DOCX 다운로드", data=final_docx, file_name=f"[최종본] {project_data['topic']}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
                with col2:
                    st.download_button(label="📥 최종본 TXT 다운로드", data=display_text.encode('utf-8'), file_name=f"[최종본] {project_data['topic']}.txt", mime="text/plain", use_container_width=True)
        else: st.error("선택된 프로젝트의 본문 내용을 찾을 수 없습니다.")
    else: st.warning("먼저 '1단계'에서 프로젝트를 시작해주세요.")

elif st.session_state.active_tab == "5단계: PPT 전환":
    st.header("📝 제안서 기반 PPT 자동 전환")

    if 'project_id' in st.session_state:
        project_id = st.session_state.project_id
        project_data, stages = get_project_data(project_id)
        proposal_content = stages.get("4단계: 최종본", stages.get("3단계: 본문 생성"))

        if project_data and proposal_content:
            st.info(f"**PPT 전환 대상:** {project_data['topic']} (ID: {project_id})")
            
            if "4단계: 최종본" in stages:
                st.success("✅ AI가 개선한 최종본을 사용하여 PPT를 생성합니다.")
            else:
                st.warning("⚠️ AI 개선을 거치지 않은 초안을 사용하여 PPT를 생성합니다. 더 높은 품질을 위해 4단계 진행을 권장합니다.")

            with st.expander("제안서 원본 보기"):
                st.text(proposal_content)

            st.subheader("PPT 설정")
            ppt_page_count = st.number_input("목표 슬라이드 수", min_value=5, max_value=20, value=10)
            theme_name = st.selectbox("디자인 테마 선택", options=["심플 (파랑)", "비즈니스 (회색)", "크리에이티브 (보라)", "자연 (초록)", "따뜻함 (주황)"])
            
            if st.button("🚀 PPT 생성 시작하기", type="primary", use_container_width=True):
                if 'generated_ppt_file' in st.session_state:
                    del st.session_state['generated_ppt_file']
                with st.status("PPT 생성을 시작합니다...", expanded=True) as status:
                    status.write("1/3: AI가 제안서 내용을 분석 및 요약 중입니다...")
                    slides_json_str = generate_ppt_slides(proposal_content, ppt_page_count)
                    if slides_json_str:
                        status.write("2/3: 요약된 내용을 바탕으로 PPT 파일을 조립하고 있습니다...")
                        ppt_file = create_ppt_presentation(slides_json_str, theme_name)
                        if ppt_file:
                            st.session_state.generated_ppt_file = ppt_file
                            status.write("3/3: 생성 완료! 다운로드 버튼을 준비합니다.")
                            status.update(label="🎉 PPT 생성이 완료되었습니다!", state="complete")
                        else:
                            status.update(label="오류: PPT 파일 조립 실패", state="error")
                    else:
                        status.update(label="오류: AI 요약 실패", state="error")

            if 'generated_ppt_file' in st.session_state and st.session_state.generated_ppt_file:
                st.success("이제 아래에서 PPT 파일을 다운로드할 수 있습니다.")
                
                safe_topic = re.sub(r'[^\w\d_.-]', '_', project_data['topic'])
                st.download_button(
                    label="📥 PPTX 파일 다운로드",
                    data=st.session_state.generated_ppt_file,
                    file_name=f"{safe_topic}_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
        else:
            st.error("선택된 프로젝트의 본문 내용을 찾을 수 없습니다.")
    else:
        st.warning("먼저 '1단계'에서 프로젝트를 시작하거나, History에서 PPT로 전환할 프로젝트를 선택해주세요.")